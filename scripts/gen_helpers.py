from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Constants
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
TEAL = RGBColor(0x22, 0x62, 0x7E)
RED = RGBColor(0xEC, 0x00, 0x00)
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
FONT_BODY = "Calibri"
FONT_CODE = "Consolas"
FONT_TITLE = "Calibri"
GREEN = RGBColor(0x0D, 0x94, 0x48)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
ORANGE = RGBColor(0xB4, 0x5D, 0x09)
CODE_BG = RGBColor(0x1E, 0x1E, 0x1E)
CODE_FG = RGBColor(0xD4, 0xD4, 0xD4)

OUT_DIR = "/mnt/d/Dev/Fabric/docs/slides/Modulo 4"

def new_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def _add_shape(slide, left, top, width, height, fill=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    return shape

def _add_text_box(slide, left, top, width, height, text, font_size=Pt(20), bold=False, color=BLACK, font_name=FONT_BODY, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_title_slide(prs, title, subtitle="", module="Modulo 4: Tokens y Smart Contracts"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.15), fill=TEAL)
    _add_shape(slide, Inches(0), Inches(7.35), SLIDE_W, Inches(0.15), fill=TEAL)
    _add_text_box(slide, Inches(0.5), Inches(0.4), Inches(12), Inches(0.5),
                  module, font_size=Pt(16), color=TEAL, bold=True)
    _add_text_box(slide, Inches(0.5), Inches(2.0), Inches(12), Inches(2.5),
                  title, font_size=Pt(54), bold=True, color=TEAL, alignment=PP_ALIGN.LEFT)
    if subtitle:
        _add_text_box(slide, Inches(0.5), Inches(4.5), Inches(12), Inches(1.5),
                      subtitle, font_size=Pt(24), color=DARK_GRAY)
    return slide

def add_section_slide(prs, section_title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_shape(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=TEAL)
    _add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(2),
                  section_title, font_size=Pt(48), bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)
    if subtitle:
        _add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(1.5),
                      subtitle, font_size=Pt(24), color=WHITE)
    return slide

def add_content_slide(prs, title, bullets, subbullets=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), fill=TEAL)
    _add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                  title, font_size=Pt(32), bold=True, color=TEAL)
    _add_shape(slide, Inches(0.5), Inches(1.15), Inches(12), Inches(0.03), fill=TEAL)
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(11.5), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(20)
        p.font.name = FONT_BODY
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(8)
        p.level = 0
        if subbullets and i in subbullets:
            for sub in subbullets[i]:
                sp = tf.add_paragraph()
                sp.text = "    " + sub
                sp.font.size = Pt(17)
                sp.font.name = FONT_BODY
                sp.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
                sp.space_after = Pt(4)
                sp.level = 1
    return slide

def add_two_column_slide(prs, title, left_title, left_bullets, right_title, right_bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), fill=TEAL)
    _add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                  title, font_size=Pt(32), bold=True, color=TEAL)
    _add_shape(slide, Inches(0.5), Inches(1.15), Inches(12), Inches(0.03), fill=TEAL)
    _add_text_box(slide, Inches(0.7), Inches(1.4), Inches(5.5), Inches(0.5),
                  left_title, font_size=Pt(22), bold=True, color=TEAL)
    txL = slide.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(5.5), Inches(5.0))
    tfL = txL.text_frame
    tfL.word_wrap = True
    for i, b in enumerate(left_bullets):
        p = tfL.paragraphs[0] if i == 0 else tfL.add_paragraph()
        p.text = b
        p.font.size = Pt(18)
        p.font.name = FONT_BODY
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(6)
    _add_shape(slide, Inches(6.5), Inches(1.4), Inches(0.03), Inches(5.5), fill=RGBColor(0xCC,0xCC,0xCC))
    _add_text_box(slide, Inches(6.9), Inches(1.4), Inches(5.5), Inches(0.5),
                  right_title, font_size=Pt(22), bold=True, color=TEAL)
    txR = slide.shapes.add_textbox(Inches(6.9), Inches(2.0), Inches(5.5), Inches(5.0))
    tfR = txR.text_frame
    tfR.word_wrap = True
    for i, b in enumerate(right_bullets):
        p = tfR.paragraphs[0] if i == 0 else tfR.add_paragraph()
        p.text = b
        p.font.size = Pt(18)
        p.font.name = FONT_BODY
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(6)
    return slide

def add_table_slide(prs, title, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), fill=TEAL)
    _add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                  title, font_size=Pt(32), bold=True, color=TEAL)
    _add_shape(slide, Inches(0.5), Inches(1.15), Inches(12), Inches(0.03), fill=TEAL)
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_w = Inches(11.5)
    row_h = Inches(0.45)
    table_h = row_h * n_rows
    top = Inches(1.5)
    left = Inches(0.9)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, table_w, table_h)
    table = table_shape.table
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = TEAL
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = FONT_BODY
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i+1, j)
            cell.text = val
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(15)
                p.font.name = FONT_BODY
                p.font.color.rgb = DARK_GRAY
    return slide

def add_code_slide(prs, title, code, language="go", subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), fill=TEAL)
    _add_text_box(slide, Inches(0.5), Inches(0.3), Inches(10), Inches(0.8),
                  title, font_size=Pt(28), bold=True, color=TEAL)
    if subtitle:
        _add_text_box(slide, Inches(0.5), Inches(0.95), Inches(10), Inches(0.4),
                      subtitle, font_size=Pt(16), color=DARK_GRAY)
    _add_shape(slide, Inches(11.5), Inches(0.35), Inches(1.3), Inches(0.4), fill=TEAL)
    _add_text_box(slide, Inches(11.5), Inches(0.35), Inches(1.3), Inches(0.4),
                  language.upper(), font_size=Pt(14), bold=True, color=WHITE,
                  alignment=PP_ALIGN.CENTER)
    code_top = Inches(1.35) if not subtitle else Inches(1.55)
    _add_shape(slide, Inches(0.5), code_top, Inches(12.3), Inches(5.6), fill=CODE_BG)
    _add_text_box(slide, Inches(0.8), code_top + Inches(0.2), Inches(11.7), Inches(5.2),
                  code, font_size=Pt(15), color=CODE_FG, font_name=FONT_CODE)
    return slide

def add_prompt_slide(prs, title, prompt_text, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), fill=TEAL)
    _add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                  title, font_size=Pt(28), bold=True, color=TEAL)
    _add_shape(slide, Inches(0.5), Inches(1.15), Inches(12), Inches(0.03), fill=TEAL)
    _add_shape(slide, Inches(0.5), Inches(1.5), Inches(2.2), Inches(0.45), fill=PURPLE)
    _add_text_box(slide, Inches(0.5), Inches(1.5), Inches(2.2), Inches(0.45),
                  "PROMPT IA", font_size=Pt(16), bold=True, color=WHITE,
                  alignment=PP_ALIGN.CENTER)
    _add_shape(slide, Inches(0.5), Inches(2.2), Inches(12.3), Inches(4.5), fill=RGBColor(0xF8, 0xF4, 0xFF))
    _add_text_box(slide, Inches(0.8), Inches(2.4), Inches(11.7), Inches(4.1),
                  prompt_text, font_size=Pt(17), color=DARK_GRAY, font_name=FONT_CODE)
    if notes:
        _add_text_box(slide, Inches(0.5), Inches(6.8), Inches(12), Inches(0.5),
                      notes, font_size=Pt(14), color=RGBColor(0x99, 0x99, 0x99))
    return slide

def add_image_placeholder(prs, title, description, prompt_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), fill=TEAL)
    _add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                  title, font_size=Pt(32), bold=True, color=TEAL)
    _add_shape(slide, Inches(0.5), Inches(1.15), Inches(12), Inches(0.03), fill=TEAL)
    _add_shape(slide, Inches(1.5), Inches(1.5), Inches(10.3), Inches(4.5), fill=LIGHT_GRAY)
    _add_text_box(slide, Inches(2), Inches(2.5), Inches(9.3), Inches(1.5),
                  description, font_size=Pt(24), color=RGBColor(0x99, 0x99, 0x99),
                  alignment=PP_ALIGN.CENTER)
    _add_shape(slide, Inches(1.5), Inches(4.2), Inches(10.3), Inches(0.35), fill=PURPLE)
    _add_text_box(slide, Inches(1.5), Inches(4.2), Inches(10.3), Inches(0.35),
                  "Prompt para IA generativa:", font_size=Pt(13), bold=True, color=WHITE,
                  alignment=PP_ALIGN.CENTER)
    _add_text_box(slide, Inches(1.8), Inches(4.7), Inches(9.7), Inches(1.0),
                  prompt_text, font_size=Pt(14), color=PURPLE, font_name=FONT_CODE)
    return slide

def add_debate_slide(prs, title, questions):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), fill=RED)
    _add_shape(slide, Inches(0.5), Inches(0.3), Inches(1.8), Inches(0.5), fill=RED)
    _add_text_box(slide, Inches(0.5), Inches(0.3), Inches(1.8), Inches(0.5),
                  "DEBATE", font_size=Pt(18), bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    _add_text_box(slide, Inches(2.6), Inches(0.3), Inches(10), Inches(0.8),
                  title, font_size=Pt(30), bold=True, color=DARK_GRAY)
    _add_shape(slide, Inches(0.5), Inches(1.15), Inches(12), Inches(0.03), fill=RED)
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(11.5), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, q in enumerate(questions):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = q
        p.font.size = Pt(22)
        p.font.name = FONT_BODY
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(16)
    return slide

def add_activity_slide(prs, title, instructions, badge_text="PRACTICA"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), fill=GREEN)
    _add_shape(slide, Inches(0.5), Inches(0.3), Inches(2.5), Inches(0.5), fill=GREEN)
    _add_text_box(slide, Inches(0.5), Inches(0.3), Inches(2.5), Inches(0.5),
                  badge_text, font_size=Pt(18), bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    _add_text_box(slide, Inches(3.3), Inches(0.3), Inches(9.5), Inches(0.8),
                  title, font_size=Pt(30), bold=True, color=DARK_GRAY)
    _add_shape(slide, Inches(0.5), Inches(1.15), Inches(12), Inches(0.03), fill=GREEN)
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(11.5), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, inst in enumerate(instructions):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = inst
        p.font.size = Pt(20)
        p.font.name = FONT_BODY
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(10)
    return slide

def add_special_slide(prs, badge_text, title, content_lines, badge_color=ORANGE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), fill=badge_color)
    _add_shape(slide, Inches(0.5), Inches(0.3), Inches(3.5), Inches(0.5), fill=badge_color)
    _add_text_box(slide, Inches(0.5), Inches(0.3), Inches(3.5), Inches(0.5),
                  badge_text, font_size=Pt(18), bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    _add_text_box(slide, Inches(4.3), Inches(0.25), Inches(8.5), Inches(0.8),
                  title, font_size=Pt(28), bold=True, color=DARK_GRAY)
    _add_shape(slide, Inches(0.5), Inches(1.15), Inches(12), Inches(0.03), fill=badge_color)
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(11.5), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(content_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(20)
        p.font.name = FONT_BODY
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(10)
    return slide

def add_review_slide(prs, title, questions):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_shape(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=TEAL)
    _add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                  title, font_size=Pt(36), bold=True, color=WHITE)
    _add_shape(slide, Inches(0.5), Inches(1.15), Inches(12), Inches(0.03), fill=WHITE)
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(11.5), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, q in enumerate(questions):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{i+1}. {q}"
        p.font.size = Pt(20)
        p.font.name = FONT_BODY
        p.font.color.rgb = WHITE
        p.space_after = Pt(10)
    return slide
