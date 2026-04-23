from pptx import Presentation

prs = Presentation("/mnt/d/Dev/Fabric/docs/slides/Modulo 3/adopcion.pptx")
total = len(prs.slides)
print(f"Total slides: {total}\n")

# Mostrar todas las slides con titulo, para ver donde estan los debates/preguntas
for idx, slide in enumerate(prs.slides):
    title = ""
    all_text = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    if not title and len(t) > 5:
                        title = t[:80]
                    all_text.append(t)
    print(f"Slide {idx+1}: {title}")

print("\n\n### DETALLE DE LAS SLIDES CON PREGUNTAS/DEBATE ###\n")
for idx, slide in enumerate(prs.slides):
    slide_text = ""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                slide_text += para.text + " "
    if "DEBATE" in slide_text or "Actividad" in slide_text or "?" in slide_text or "Pregunta" in slide_text or "Mision" in slide_text:
        print(f"\n{'='*60}")
        print(f"SLIDE {idx+1}")
        print('='*60)
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        print(f"  {t}")
