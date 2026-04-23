"""
Anade slides de respuestas a los dias 1-5 del Modulo 3.
Abre cada pptx existente y anade al final SIN tocar las existentes.
"""
import sys
sys.path.insert(0, "/mnt/d/Dev/Fabric/scripts")
from gen_helpers import add_content_slide
from pptx import Presentation

BASE = "/mnt/d/Dev/Fabric/docs/slides/Modulo 3"

# ============================================================
# DIA 1: debate (arbol decision) + repaso
# ============================================================
prs = Presentation(f"{BASE}/dia_1.pptx")

add_content_slide(prs, "Respuestas al debate: ¿blockchain o BD?", [
    "1. Hospital registra historiales de sus pacientes:",
    "2. Tres hospitales comparten historiales de pacientes compartidos:",
    "3. Banco registra sus transacciones internas:",
    "4. Diez bancos: liquidacion interbancaria sin intermediario:",
    "5. Empresa rastrea envios internos entre almacenes:",
    "6. Consorcio de 50 productores certifica origen:",
], subbullets={
    0: ["BD tradicional. Un solo actor, sin necesidad de compartir. Blockchain seria over-engineering."],
    1: ["Blockchain. Multiples orgs, datos compartidos, trazabilidad necesaria.",
        "Requiere diseño con Private Data para proteger datos personales (GDPR)."],
    2: ["BD tradicional. Un solo actor; si necesita auditoria, append-only DB o WORM storage."],
    3: ["Blockchain. Caso arquetipico (ver HKMA eTradeConnect). Evita intermediario y detecta fraude."],
    4: ["BD tradicional o ERP. Un solo actor — no hay problema de confianza."],
    5: ["Blockchain. Multiples orgs, desconfianza posible, certificacion verificable.",
        "Ver caso Walmart Food Trust."],
})

add_content_slide(prs, "Respuestas al repaso", [
    "1. Tres proyectos en produccion:",
    "2. Dos fracasos y su razon:",
    "3. Arbol de decision (5 preguntas):",
    "4. Problema del fundador dominante:",
    "5. Respaldo del regulador como factor de exito:",
], subbullets={
    0: ["Walmart Food Trust (trazabilidad alimentaria).",
        "HKMA eTradeConnect (trade finance bancos Hong Kong).",
        "ANZ + Scentre Group (tokenizacion de alquileres comerciales)."],
    1: ["TradeLens (Maersk+IBM) — gobernanza: fundador dominante ahuyento competidores.",
        "B3i (reaseguros) — modelo de negocio: sin ingresos, capital se agoto."],
    2: ["¿Multiples partes sin confianza mutua?",
        "¿Necesitan compartir datos?",
        "¿Inmutabilidad del registro es critica?",
        "¿Reglas deben ser verificables por todos?",
        "¿Ningun actor puede ser el 'dueno'?"],
    3: ["Cuando una empresa lanza y controla la red, los competidores no se suman.",
        "Ven un conflicto de intereses: dar datos a un rival.",
        "Sin competidores, la red no alcanza masa critica -> fracaso."],
    4: ["Impone adopcion (p.ej. HKMA obligo a bancos de Hong Kong).",
        "Da legitimidad legal y regulatoria.",
        "Elimina la incertidumbre sobre si el proyecto cumple la ley.",
        "Actua como arbitro en disputas."],
})

prs.save(f"{BASE}/dia_1.pptx")
print(f"dia_1.pptx: total {len(prs.slides)} slides")

# ============================================================
# DIA 2: solo repaso (5 preguntas)
# ============================================================
prs = Presentation(f"{BASE}/dia_2.pptx")

add_content_slide(prs, "Respuestas al repaso (1/2)", [
    "1. Cuando tiene sentido blockchain enterprise (3 condiciones):",
    "2. Cuando NO tiene sentido (3 razones):",
    "3. Como presentarlo a un comite de direccion:",
], subbullets={
    0: ["Hay multiples organizaciones que necesitan compartir datos.",
        "Existe desconfianza mutua (ninguna quiere que otra controle la BD).",
        "La inmutabilidad/auditabilidad es critica (regulacion, compliance)."],
    1: ["Un solo actor gestiona todos los datos -> BD tradicional.",
        "Los datos no se comparten entre orgs -> APIs entre sistemas.",
        "No hay problema de confianza -> SaaS compartido con permisos.",
        "(Tambien: si el rendimiento es critico >10k TPS)."],
    2: ["Empezar por EL PROBLEMA de negocio, no por la tecnologia.",
        "Cuantificar el coste actual (euros, tiempo, fraude).",
        "Por que las soluciones actuales no funcionan.",
        "Propuesta: 'registro distribuido compartido' (evita la palabra 'blockchain').",
        "Socios potenciales, ROI estimado, riesgos y mitigaciones."],
})

add_content_slide(prs, "Respuestas al repaso (2/2)", [
    "4. Metricas para medir el exito:",
    "5. ¿Por que no empezar con la palabra 'blockchain'?",
], subbullets={
    0: ["Eficiencia: tiempo de proceso reducido (dias/semanas -> segundos/horas).",
        "Errores: reduccion de errores manuales y disputas.",
        "Fraude: casos detectados/evitados antes de su impacto.",
        "Costes: intermediarios eliminados, auditorias mas baratas.",
        "Payback: tiempo para recuperar la inversion inicial."],
    1: ["'Blockchain' genera ruido: se asocia con cripto, especulacion, scams.",
        "Los directivos no tecnicos pueden cerrarse antes de entender el problema.",
        "Mejor: 'registro distribuido', 'ledger compartido', 'plataforma colaborativa'.",
        "Presenta la SOLUCION a su problema, no la tecnologia.",
        "Despues de aprobacion, explicar que la tecnologia base es blockchain/Fabric."],
})

prs.save(f"{BASE}/dia_2.pptx")
print(f"dia_2.pptx: total {len(prs.slides)} slides")

# ============================================================
# DIA 3: debate (gobernanza KYC) + repaso
# ============================================================
prs = Presentation(f"{BASE}/dia_3.pptx")

add_content_slide(prs, "Respuestas al debate: gobernanza consorcio KYC", [
    "1. Modelo de gobernanza:",
    "2. Politica de endorsement:",
    "3. ¿Quien opera los orderers?",
    "4. Como incorporar un nuevo banco:",
    "5. Que pasa si un banco quiere salir:",
], subbullets={
    0: ["Federado con comite rotativo: evita que un banco domine.",
        "Alternativa: democratico (1 banco = 1 voto) si son comparables en tamano.",
        "Fundador dominante: evitar — replicaria el problema TradeLens."],
    1: ["Lectura KYC: OR(banco_titular) — el banco del cliente basta.",
        "Actualizar KYC: AND(banco_titular, banco_validador) — doble verificacion.",
        "Cambios de politica del canal: MAJORITY de los 5 bancos."],
    2: ["Distribuidos: cada banco opera un nodo Raft (5 nodos total).",
        "Ningun banco controla el ordering solo.",
        "Alta disponibilidad: tolera 2 fallos simultaneos."],
    3: ["Evaluacion por comite (sistemas, regulacion, solvencia).",
        "Aprobacion formal: MAJORITY de los 5 bancos existentes.",
        "Onboarding tecnico: emitir certs, unir al canal, instalar chaincode.",
        "Cuota escalada: paga proporcional a su tamano."],
    4: ["Datos KYC ya registrados: inmutables, no se pueden borrar.",
        "Su certificado se revoca (CRL) — no puede leer ni escribir.",
        "Los clientes del banco migran a otro banco del consorcio o salen del sistema.",
        "Responsabilidad legal por datos antiguos queda segun el acuerdo del consorcio."],
})

add_content_slide(prs, "Respuestas al repaso", [
    "1. Gobernanza y por que es critica:",
    "2. Tres modelos de gobernanza + ejemplo:",
    "3. ¿Como codifica Fabric la gobernanza?",
    "4. ¿Quien deberia pagar los orderers?",
    "5. ¿Que definir ANTES de que una org abandone?",
], subbullets={
    0: ["Conjunto de reglas que definen como se toman decisiones en la red.",
        "Sin ella, no hay consenso sobre cambios y el consorcio se paraliza.",
        "Proyectos como TradeLens cerraron por falla de gobernanza, no tecnologia."],
    1: ["Fundador dominante: Walmart Food Trust (Walmart impone).",
        "Democratico: Alastria (cada miembro un voto en asamblea).",
        "Federado: EBSI (comite ejecutivo rotativo entre miembros)."],
    2: ["Politicas de canal (MAJORITY Admins, etc).",
        "Politicas de endorsement por chaincode o state-based por key.",
        "Lifecycle: approveformyorg + commit requieren consenso.",
        "Configuracion del canal: cambios via propuesta + firmas + update."],
    3: ["Depende del modelo: cuota fija por org o proporcional al uso.",
        "Modelo sponsor: una org paga todo al principio (no recomendado a largo plazo).",
        "Producci­on: cada org opera su propio orderer (Raft distribuido)."],
    4: ["Como se revocan sus certificados y se actualiza la CRL.",
        "Que pasa con los datos en el ledger (inmutables, no se borran).",
        "Quien asume sus responsabilidades operativas.",
        "Clausula de resolucion de disputas y sucesion.",
        "Todo esto DEBE estar en el acuerdo del consorcio ANTES de empezar."],
})

prs.save(f"{BASE}/dia_3.pptx")
print(f"dia_3.pptx: total {len(prs.slides)} slides")

# ============================================================
# DIA 4: solo repaso (5 preguntas)
# ============================================================
prs = Presentation(f"{BASE}/dia_4.pptx")

add_content_slide(prs, "Respuestas al repaso (1/2)", [
    "1. GDPR vs inmutabilidad de blockchain:",
    "2. ¿Smart contract es un contrato legal? ¿Que falta?",
    "3. ¿Que es MiCA y a quien aplica?",
], subbullets={
    0: ["No son incompatibles, pero requieren diseño cuidadoso.",
        "NUNCA guardar datos personales on-chain — solo hash o referencia.",
        "Datos reales off-chain (BD tradicional), borrables para cumplir derecho al olvido.",
        "Private Data Collections con blockToLive permiten caducar datos.",
        "Al borrar dato off-chain, el hash on-chain queda huerfano (pseudonimizacion)."],
    1: ["NO, al menos no automaticamente. Un smart contract es codigo.",
        "Para tener valor legal necesita: consentimiento explicito de las partes,",
        "identificacion (no pseudonima), objeto licito, jurisdiccion clara,",
        "mecanismo de resolucion de disputas.",
        "Enfoque practico: contrato legal tradicional + smart contract que lo ejecuta."],
    2: ["Markets in Crypto-Assets Regulation (diciembre 2024).",
        "Aplica a: tokens de utilidad, stablecoins, criptoactivos publicos, exchanges.",
        "NO aplica directamente a blockchain enterprise permissioned sin tokens publicos.",
        "Relevante si tu proyecto Fabric emite tokens que salen al mercado.",
        "Requisitos: whitepaper, proteccion al inversor, autorizacion ESMA/CNMV."],
})

add_content_slide(prs, "Respuestas al repaso (2/2)", [
    "4. Regulacion si tokenizas un bono en Fabric:",
    "5. Si un chaincode tiene un bug, ¿quien responde?",
], subbullets={
    0: ["Si el token representa un valor negociable -> MiFID II aplica.",
        "Requiere autorizacion de la CNMV (Espana) o equivalente europeo.",
        "Aplican las mismas reglas que a valores tradicionales (folleto, KYC, etc).",
        "Si es una STO (Security Token Offering): mas regulado que ICO pero con marco.",
        "Tokenizacion interna entre orgs (no ventas publicas): menos regulacion directa."],
    1: ["Responsabilidad compartida: todas las orgs que aprobaron el lifecycle.",
        "Preguntas clave a definir en el acuerdo del consorcio:",
        "  - ¿Quien audita los chaincodes antes de aprobarlos?",
        "  - ¿Hay seguro de responsabilidad civil?",
        "  - ¿Mecanismo de compensacion?",
        "  - ¿Jurisdiccion y ley aplicables?",
        "En la practica: el desarrollador/proveedor responde contractualmente.",
        "El consorcio establece limites via SLA y clausulas de indemnizacion."],
})

prs.save(f"{BASE}/dia_4.pptx")
print(f"dia_4.pptx: total {len(prs.slides)} slides")

# ============================================================
# DIA 5: debate (cierre modulo) + repaso del modulo
# ============================================================
prs = Presentation(f"{BASE}/dia_5.pptx")

add_content_slide(prs, "Respuestas al debate de cierre", [
    "1. Lo mas sorprendente del modulo:",
    "2. ¿Han cambiado de opinion sobre blockchain enterprise?",
    "3. Oportunidades en vuestro sector:",
    "4. ¿Que es mas dificil?",
    "5. Proyecto con presupuesto ilimitado:",
], subbullets={
    0: ["Normalmente: que la mayoria de fracasos NO son por tecnologia.",
        "Son por gobernanza, incentivos, regulacion o UX.",
        "La tecnologia Fabric funciona — el problema es el contexto."],
    1: ["Depende del alumno, pero patron comun:",
        "  - De 'revolucionario' a 'util en ciertos casos concretos'.",
        "  - De 'solucion a todo' a 'herramienta para problemas especificos'.",
        "  - Muchos empiezan esceptico con cripto y terminan valorando enterprise."],
    2: ["Sin respuesta unica — depende del sector de cada alumno.",
        "Preguntas guia: ¿tenemos multiples actores con desconfianza? ¿trazabilidad?",
        "Si la respuesta es si a ambas, probablemente si."],
    3: ["Consenso habitual: GOBERNANZA > REGULACION > TECNOLOGIA.",
        "La tecnologia esta madura; la gobernanza no siempre se resuelve bien.",
        "La regulacion depende de la jurisdiccion y evoluciona, no es un bloqueo absoluto."],
    4: ["Trazabilidad alimentaria a escala europea (Food Trust pan-europeo).",
        "KYC compartido entre bancos europeos.",
        "Registro de propiedad inmobiliaria tokenizado.",
        "Creditos de carbono auditables para cumplimiento climatico.",
        "Interoperabilidad entre redes: Alastria-EBSI-HKMA interoperables."],
})

add_content_slide(prs, "Respuestas al repaso del Modulo 3 (1/2)", [
    "1. Casos de exito/fracaso — patrones comunes:",
    "2. Arbol de decision ¿blockchain o BD?",
    "3. Como presentar a direccion:",
    "4. Modelos de gobernanza:",
], subbullets={
    0: ["Los que funcionan: consorcio con incentivos alineados, regulacion compatible,",
        "UX transparente, MVP pequeno, gobernanza equilibrada.",
        "Los que fracasan: fundador dominante, sin modelo de ingresos,",
        "regulacion incompatible, UX compleja, intentar cubrir todo el mercado."],
    1: ["5 preguntas: partes sin confianza / compartir datos / inmutabilidad critica /",
        "reglas verificables / ningun dueno unico. Si a todas -> blockchain."],
    2: ["Empezar por el problema y el coste actual.",
        "Evitar la palabra 'blockchain' al inicio.",
        "Hablar de 'registro distribuido compartido'.",
        "Presentar ROI, socios, MVP, riesgos."],
    3: ["Fundador dominante: Walmart Food Trust.",
        "Democratico: Alastria (1 miembro = 1 voto).",
        "Federado: EBSI (comite ejecutivo rotativo)."],
})

add_content_slide(prs, "Respuestas al repaso del Modulo 3 (2/2)", [
    "5. GDPR vs inmutabilidad:",
    "6. MiCA y regulacion de tokens:",
    "7. Metodologia: de la idea al roadmap:",
    "8. Costes tipicos de un proyecto Fabric:",
], subbullets={
    0: ["Datos personales off-chain. Hash on-chain.",
        "Private Data con blockToLive caduca datos privados.",
        "Borrar off-chain deja hash huerfano = pseudonimizacion efectiva."],
    1: ["MiCA aplica a cripto publica (stablecoins, exchanges, tokens).",
        "Fabric enterprise SIN tokens publicos: menos regulacion directa.",
        "Tokenizar valores (bonos, acciones) en Fabric: MiFID II + autorizacion CNMV."],
    2: ["Descubrimiento (1-2 meses): problema, validar si blockchain aporta, socios.",
        "PoC (3-4 meses): red minima, chaincode basico, demostrar funciona.",
        "MVP (5-8 meses): red con Fabric CA, chaincode completo, app, tests.",
        "Piloto (9-12 meses): datos reales, mas orgs, integracion legacy.",
        "Produccion (12+ meses): SLAs, monitorizacion, soporte 24/7."],
    3: ["Infraestructura: 500-2000€/mes por org (cloud), menos on-premise.",
        "Desarrollo: chaincode 2-6 meses + app cliente 1-3 meses.",
        "Integracion con legacy: el coste mas impredecible (depende del sistema).",
        "Operaciones: 0.5-1 persona por org + 20-30% del dev/ano mantenimiento.",
        "Gobernanza: reuniones, comites, decision-making (tiempo directivos)."],
})

prs.save(f"{BASE}/dia_5.pptx")
print(f"dia_5.pptx: total {len(prs.slides)} slides")
