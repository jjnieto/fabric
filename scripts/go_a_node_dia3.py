"""
Convierte el codigo Go a Node.js en las slides del dia 3 del Modulo 4.
Abre el pptx existente y cambia SOLO el texto de los shapes identificados,
preservando formato (fuente, tamano, color, posicion, tamano del shape).
"""
from pptx import Presentation
from copy import deepcopy

PPTX_PATH = "/mnt/d/Dev/Fabric/docs/slides/Modulo 4/dia_3.pptx"

# Mapping: (slide_index, shape_index) -> texto nuevo (0-based)
REPLACEMENTS = {
    # ============================================================
    # Slide 6 (idx 5): Control de acceso por organizacion
    # ============================================================
    (5, 4): "NODE.JS",
    (5, 6): """async TransferProperty(ctx, id, newOwner) {
    // Leer propiedad
    const property = await this.ReadProperty(ctx, id);

    // Verificar: solo el propietario actual puede transferir
    const clientMSP = ctx.clientIdentity.getMSPID();
    if (property.owner !== clientMSP) {
        throw new Error(
            `solo el propietario actual (${property.owner}) ` +
            `puede transferir`);
    }""",
    (5, 7): """    // Verificar: la propiedad debe estar en estado 'active'
    if (property.status !== 'active') {
        throw new Error(
            `la propiedad esta en estado '${property.status}', ` +
            `no se puede transferir`);
    }

    // Transferir
    property.owner = newOwner;
    property.status = 'active';
    await ctx.stub.putState(id,
        Buffer.from(JSON.stringify(property)));
}""",

    # ============================================================
    # Slide 7 (idx 6): Control de acceso ABAC
    # ============================================================
    (6, 4): "NODE.JS",
    (6, 6): """// Solo un usuario con rol "registrador" puede registrar propiedades
async RegisterProperty(ctx, id, address, owner) {
    // Verificar atributo "role" en el certificado X.509
    const ok = ctx.clientIdentity
        .assertAttributeValue('role', 'registrador');
    if (!ok) {
        throw new Error(
            'solo registradores pueden crear propiedades');
    }

    // El atributo se anade al certificado durante el enrollment con Fabric CA:
    //   fabric-ca-client enroll ... --enrollment.attrs "role=registrador"

    const property = {
        docType: 'property',
        id,
        address,
        owner,
        status: 'active'
    };

    await ctx.stub.putState(id,
        Buffer.from(JSON.stringify(property)));
}""",

    # ============================================================
    # Slide 10 (idx 9): Maquina de estados en codigo
    # ============================================================
    (9, 3): "NODE.JS",
    (9, 5): """// Mapa de transiciones validas
const validTransitions = {
    active:       ['transferring', 'blocked'],
    transferring: ['active'],
    blocked:      ['active'],
};

function isValidTransition(currentStatus, newStatus) {
    const allowed = validTransitions[currentStatus];
    if (!allowed) {
        return false;
    }
    return allowed.includes(newStatus);
}""",
    (9, 6): """// Uso en una funcion del chaincode:
if (!isValidTransition(property.status, newStatus)) {
    throw new Error(
        `transicion no permitida: ` +
        `${property.status} -> ${newStatus}`);
}""",

    # ============================================================
    # Slide 14 (idx 13): Private Data leer y escribir
    # ============================================================
    (13, 4): "NODE.JS",
    (13, 6): """// Escribir en una Private Data Collection
async SetPrivatePrice(ctx, propertyID) {
    // Los datos privados vienen en transient data (no en el ledger)
    const transientMap = ctx.stub.getTransient();
    const priceData = transientMap.get('price');
    if (!priceData) {
        throw new Error(
            'falta el campo "price" en transient data');
    }

    // Guardar en la coleccion privada
    await ctx.stub.putPrivateData(
        'priceAgreement', propertyID, priceData);
}""",
    (13, 7): """// Leer de una Private Data Collection
async GetPrivatePrice(ctx, propertyID) {
    const priceData = await ctx.stub.getPrivateData(
        'priceAgreement', propertyID);

    if (!priceData || priceData.length === 0) {
        return '';
    }
    return priceData.toString();
}""",

    # ============================================================
    # Slide 16 (idx 15): Eventos ejemplo completo
    # ============================================================
    (15, 3): "NODE.JS",
    (15, 5): """// En el chaincode: emitir evento al transferir
// (en JS no hay struct, se construye el objeto directamente)

async TransferProperty(ctx, id, newOwner) {
    // ... logica de transferencia ...

    const event = {
        propertyID: id,
        from: property.owner,
        to: newOwner,
        timestamp: new Date().toISOString()
    };""",
    (15, 6): """    const eventJSON = Buffer.from(JSON.stringify(event));
    ctx.stub.setEvent('PropertyTransferred', eventJSON);
}

// En el cliente Node.js: escuchar eventos
// const listener = async (event) => {
//     const payload = JSON.parse(
//         event.payload.toString());
//     console.log('Transferencia:', payload.propertyID);
// };
// contract.addContractListener(listener);""",
}


def replace_text_preserving_format(text_frame, new_text):
    """Reemplaza texto preservando formato del primer run."""
    first_run_format = None
    first_para_align = None
    for para in text_frame.paragraphs:
        if para.runs:
            first_run = para.runs[0]
            first_para_align = para.alignment
            rPr = first_run._r.find('{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
            if rPr is not None:
                first_run_format = deepcopy(rPr)
            break

    first_pPr = None
    for para in text_frame.paragraphs:
        pPr = para._pPr
        if pPr is not None:
            first_pPr = deepcopy(pPr)
            break

    txBody = text_frame._txBody
    for p in txBody.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}p'):
        txBody.remove(p)

    lines = new_text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
        else:
            p = text_frame.add_paragraph()

        if first_para_align is not None:
            p.alignment = first_para_align

        if first_pPr is not None:
            existing_pPr = p._pPr
            if existing_pPr is not None:
                p._p.remove(existing_pPr)
            p._p.insert(0, deepcopy(first_pPr))

        run = p.add_run()
        run.text = line
        if first_run_format is not None:
            existing_rPr = run._r.find('{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
            if existing_rPr is not None:
                run._r.remove(existing_rPr)
            run._r.insert(0, deepcopy(first_run_format))


def main():
    prs = Presentation(PPTX_PATH)

    for (slide_idx, shape_idx), new_text in REPLACEMENTS.items():
        try:
            slide = prs.slides[slide_idx]
            shape = slide.shapes[shape_idx]
            if not shape.has_text_frame:
                print(f"  ! Slide {slide_idx+1} shape {shape_idx}: no tiene text_frame, saltando")
                continue
            replace_text_preserving_format(shape.text_frame, new_text)
            preview = new_text.split('\n')[0][:50]
            print(f"  OK Slide {slide_idx+1} shape {shape_idx}: '{preview}...'")
        except IndexError:
            print(f"  ! Slide {slide_idx+1} shape {shape_idx}: no existe")

    prs.save(PPTX_PATH)
    print(f"\nGuardado: {PPTX_PATH}")


if __name__ == "__main__":
    main()
