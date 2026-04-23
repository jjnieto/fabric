from pptx import Presentation

prs = Presentation("/mnt/d/Dev/Fabric/docs/slides/Modulo 4/dia_4.pptx")
targets = [6, 7, 11, 16, 17, 18]

for idx, slide in enumerate(prs.slides):
    if (idx+1) not in targets:
        continue
    print("=" * 60)
    print(f"SLIDE {idx+1}")
    print("=" * 60)
    for s_idx, shape in enumerate(slide.shapes):
        if shape.has_text_frame:
            paras = [p.text for p in shape.text_frame.paragraphs]
            combined = "\n".join(paras)
            if len(combined.strip()) > 100 and ("func " in combined or "type " in combined):
                print(f"\n--- Shape {s_idx} ({shape.width/914400:.1f}x{shape.height/914400:.1f}) ---")
                print(combined)
    print()
