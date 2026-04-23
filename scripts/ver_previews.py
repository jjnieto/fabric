from pptx import Presentation
import os

previews = [
    "preview01_walmart.pptx",
    "preview02_wetrade.pptx",
    "preview03_hkma.pptx",
    "preview04_tradelens.pptx",
    "preview05_b3i.pptx",
]

for fname in previews:
    path = f"/mnt/d/Dev/Fabric/docs/slides/Modulo 3/{fname}"
    prs = Presentation(path)
    print("\n" + "=" * 70)
    print(f"{fname}  (total: {len(prs.slides)} slides)")
    print("=" * 70)

    for idx, slide in enumerate(prs.slides):
        slide_text = ""
        title = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        if not title and len(t) > 5:
                            title = t[:80]
                        slide_text += t + " | "
        # Solo mostrar slides con preguntas/debate/actividad
        if "?" in slide_text or "DEBATE" in slide_text or "ACTIVIDAD" in slide_text or "Pregunta" in slide_text or "Mision" in slide_text or "EJERCICIO" in slide_text or "TREASURE" in slide_text:
            print(f"\n--- Slide {idx+1}: {title} ---")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            print(f"  {t}")
