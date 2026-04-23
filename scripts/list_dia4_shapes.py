from pptx import Presentation

prs = Presentation("/mnt/d/Dev/Fabric/docs/slides/Modulo 4/dia_4.pptx")
targets = [6, 7, 11, 16, 17, 18]

for idx, slide in enumerate(prs.slides):
    if (idx+1) not in targets:
        continue
    print(f"\n=== SLIDE {idx+1} ===")
    for s_idx, shape in enumerate(slide.shapes):
        if shape.has_text_frame:
            lines = [p.text for p in shape.text_frame.paragraphs if p.text.strip()]
            first = lines[0][:60] if lines else ""
            try:
                l = shape.left/914400
                t = shape.top/914400
                w = shape.width/914400
                h = shape.height/914400
            except:
                l=t=w=h=0
            print(f"  Shape {s_idx}: pos=({l:.1f},{t:.1f}) size=({w:.1f}x{h:.1f}) lines={len(lines)} first=\"{first}\"")
