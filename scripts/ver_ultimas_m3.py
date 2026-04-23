from pptx import Presentation

for fname in ["dia_1.pptx", "dia_2.pptx", "dia_3.pptx", "dia_4.pptx", "dia_5.pptx"]:
    path = f"/mnt/d/Dev/Fabric/docs/slides/Modulo 3/{fname}"
    prs = Presentation(path)
    total = len(prs.slides)
    print("\n" + "=" * 70)
    print(f"{fname}  (total: {total} slides)")
    print("=" * 70)

    for idx in [total-2, total-1]:
        slide = prs.slides[idx]
        print(f"\n--- Slide {idx+1} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        print(f"  {t}")
