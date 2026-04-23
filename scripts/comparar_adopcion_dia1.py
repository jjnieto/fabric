from pptx import Presentation

def listar(path, nombre):
    prs = Presentation(path)
    print(f"\n{'=' * 70}")
    print(f"{nombre} ({len(prs.slides)} slides)")
    print('=' * 70)
    for idx, slide in enumerate(prs.slides):
        title = ""
        all_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        all_text.append(t)
                        if not title and len(t) > 5:
                            title = t[:70]
        # Contar imagenes
        n_pics = sum(1 for s in slide.shapes if s.shape_type == 13)  # 13 = PICTURE
        tag = f" [IMG:{n_pics}]" if n_pics else ""
        print(f"  {idx+1:2d}. {title}{tag}")

listar("/mnt/d/Dev/Fabric/docs/slides/Modulo 3/adopcion.pptx", "adopcion.pptx")
listar("/mnt/d/Dev/Fabric/docs/slides/Modulo 3/dia_1.pptx", "dia_1.pptx")
