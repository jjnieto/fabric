"""
Anade slides al final de pki_certificados_msp.pptx explicando la diferencia
entre cert intermedio y cert end-entity.
ABRE el archivo existente, NO lo regenera. NO toca slides existentes.
"""
import sys
sys.path.insert(0, "/mnt/d/Dev/Fabric/scripts")
from gen_helpers import add_content_slide, add_table_slide, add_section_slide
from pptx import Presentation

PPTX_PATH = "/mnt/d/Dev/Fabric/docs/slides/Modulo 5/pki_certificados_msp.pptx"

# Abrir archivo existente — NO crear uno nuevo
prs = Presentation(PPTX_PATH)
slides_antes = len(prs.slides)

# Slide 1: introducción al tema
add_section_slide(prs, "Anexo: Cert intermedio\nvs Cert end-entity",
    "La diferencia clave: ¿quién puede firmar otros certs?")

# Slide 2: el campo Basic Constraints
add_content_slide(prs, "El campo clave: Basic Constraints", [
    "Todos los certificados X.509 tienen un campo llamado 'Basic Constraints'",
    "Ese campo indica si el cert puede firmar OTROS certificados o no",
    "",
    "CA: TRUE → este cert puede firmar otros (es una CA)",
    "CA: FALSE → este cert NO puede firmar otros, solo se usa para si mismo",
    "",
    "Es lo que distingue funcionalmente a un cert intermedio (CA: TRUE)",
    "de un cert end-entity como el de un peer o usuario (CA: FALSE).",
])

# Slide 3: tabla comparativa
add_table_slide(prs, "Tabla comparativa de los 3 tipos de cert",
    ["Aspecto", "Cert raíz", "Cert intermedio", "End-entity"],
    [
        ["¿Auto-firmado?", "Si (Subject = Issuer)", "No (firmado por raíz)", "No (firmado por CA)"],
        ["Basic Constraints CA", "TRUE", "TRUE", "FALSE"],
        ["¿Puede firmar otros certs?", "Si", "Si", "NO"],
        ["¿Donde se almacena?", "Offline, en HSM/boveda", "Online, en servidor CA", "En el nodo (peer, app)"],
        ["Validez tipica", "10-20 años", "5-10 años", "1-2 años"],
        ["Para qué se usa", "Anclar la confianza", "Emitir end-entities", "Firmar transacciones / TLS"],
    ])

# Slide 4: cert intermedio vs end-entity (la pregunta del alumno)
add_content_slide(prs, "Intermedio vs end-entity: la diferencia real", [
    "Aspecto comun:",
    "Aspecto que los diferencia:",
    "Si un end-entity intentara emitir un cert:",
], subbullets={
    0: ["Ambos son firmados por un cert superior",
        "Ambos son archivos .pem identicos en formato",
        "Ambos tienen Subject, Issuer, clave publica, validez..."],
    1: ["El intermedio tiene 'CA: TRUE' y la opcion 'keyCertSign' en keyUsage",
        "El end-entity tiene 'CA: FALSE' y NO puede firmar otros certs",
        "Es UN solo campo en el cert lo que cambia su comportamiento"],
    2: ["1. End-entity (Alice) firma un cert para Bob",
        "2. Bob presenta su cert a Carlos para validarlo",
        "3. Carlos recorre la cadena: Bob -> firmado por Alice",
        "4. Carlos comprueba el cert de Alice: CA: FALSE",
        "5. RECHAZADO — un end-entity no puede ser CA",
        "Por eso es seguro distribuir certs end-entity sin riesgo de cadena ilegitima"],
})

# Slide 5: practica con openssl
add_content_slide(prs, "Como verlo en la practica con openssl", [
    "Inspecciona un cert de un peer (end-entity) en la red FidelityChain:",
    "",
    "openssl x509 -in cert.pem -text -noout | grep -A 2 \"Basic Constraints\"",
    "",
    "Resultado para un peer (end-entity):",
    "    X509v3 Basic Constraints: critical",
    "        CA:FALSE",
    "",
    "Compara con la CA raíz de la org:",
    "    X509v3 Basic Constraints: critical",
    "        CA:TRUE",
    "",
    "Esa unica linea (CA:TRUE vs CA:FALSE) marca toda la diferencia.",
    "Es lo que protege la cadena de confianza de delegaciones no autorizadas.",
])

# Slide 6: en Fabric con cryptogen vs Fabric CA
add_content_slide(prs, "En Fabric: ¿cuando ves cada tipo?", [
    "Con cryptogen (desarrollo, aula):",
    "Con Fabric CA en produccion seria:",
    "Como verificar si tu cert es CA o end-entity:",
], subbullets={
    0: ["1 cert raíz por org → CA: TRUE (en cacerts/)",
        "Certs de peers, admins, users → CA: FALSE (end-entity)",
        "Solo 2 niveles, sin intermedios"],
    1: ["1 cert raíz por org → CA: TRUE, offline en HSM",
        "1+ certs intermedios → CA: TRUE, en servidor CA",
        "Certs de peers, admins, users → CA: FALSE",
        "3+ niveles, con intermedios protegiendo la raíz"],
    2: ["openssl x509 -in cert.pem -text -noout | grep -i 'CA:'",
        "Si dice CA: TRUE → es una CA (raíz o intermedio)",
        "Si dice CA: FALSE → es un end-entity"],
})

prs.save(PPTX_PATH)
slides_despues = len(prs.slides)
print(f"Anadidas {slides_despues - slides_antes} slides al final de {PPTX_PATH}")
print(f"Total slides: {slides_despues} (antes: {slides_antes})")
