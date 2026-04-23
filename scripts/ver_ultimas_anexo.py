from pptx import Presentation

prs = Presentation("/mnt/d/Dev/Fabric/docs/slides/Modulo 4/dia1_anexo.pptx")
total = len(prs.slides)
print(f"Total slides: {total}\n")

for idx in [total-2, total-1]:
    slide = prs.slides[idx]
    print("=" * 60)
    print(f"SLIDE {idx+1}")
    print("=" * 60)
    for s_idx, shape in enumerate(slide.shapes):
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    print(f"  {t}")
    print()
