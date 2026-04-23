from pptx import Presentation

prs = Presentation("/mnt/d/Dev/Fabric/docs/slides/Modulo 4/dia_4.pptx")
# Para cada slide, extraer shapes 5, 6 y 7 (los que pueden tener codigo)
targets = [6, 7, 11, 16, 17, 18]

for idx, slide in enumerate(prs.slides):
    if (idx+1) not in targets:
        continue
    print("=" * 60)
    print(f"SLIDE {idx+1}")
    print("=" * 60)
    for s_idx in [5, 6, 7]:
        try:
            shape = slide.shapes[s_idx]
            if shape.has_text_frame:
                paras = [p.text for p in shape.text_frame.paragraphs]
                combined = "\n".join(paras)
                if len(combined.strip()) > 30:
                    print(f"\n--- Shape {s_idx} ---")
                    print(combined)
        except IndexError:
            pass
    print()
