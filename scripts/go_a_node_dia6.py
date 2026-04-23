"""
Convierte Go a Node.js en slides del dia 6 del Modulo 4.
Incluye cambio de titulo ("en Go" -> "en Node.js" / "(Go)" -> "(Node.js)").
"""
from pptx import Presentation
from copy import deepcopy

PPTX_PATH = "/mnt/d/Dev/Fabric/docs/slides/Modulo 4/dia_6.pptx"

REPLACEMENTS = {
    # Slide 8 (idx 7): Gateway SDK: conexion en Go
    (7, 1): "Gateway SDK: conexion en Node.js",  # titulo
    (7, 4): "NODE.JS",  # badge
    (7, 6): """const { connect, signers } = require(
    '@hyperledger/fabric-gateway');
const grpc = require('@grpc/grpc-js');
const crypto = require('crypto');

async function main() {
    // 1. Conexion gRPC al peer
    const client = new grpc.Client(
        'localhost:7051',
        grpc.credentials.createSsl(tlsRootCert));

    // 2. Crear identidad
    const identity = {
        mspId: 'Org1MSP',
        credentials: certificate
    };
    const signer = signers.newPrivateKeySigner(privateKey);""",
    (7, 7): """    // 3. Conectar al Gateway
    const gateway = connect({ client, identity, signer });

    // 4. Obtener canal y contrato
    const network = gateway.getNetwork('mychannel');
    const contract = network.getContract('foodtrace');

    // 5. Evaluar (lectura)
    const result = await contract.evaluateTransaction(
        'GetAllLots');
    console.log(new TextDecoder().decode(result));

    // 6. Submit (escritura)
    await contract.submitTransaction('ProduceLot',
        'LOT-001', 'mango', 'Spain', '500.0');
}

main().catch(console.error);""",

    # Slide 13 (idx 12): Unit test con mock del Stub (Go)
    (12, 1): "Unit test con mock del Stub (Node.js)",  # titulo
    (12, 4): "NODE.JS",  # badge
    (12, 6): """const { expect } = require('chai');
const sinon = require('sinon');
const { SmartContract } = require('../lib/smartContract');

describe('CreateProperty', () => {
    let ctx, stub, identity;

    beforeEach(() => {
        stub = {
            getState: sinon.stub(),
            putState: sinon.stub(),
        };
        identity = {
            getMSPID: sinon.stub().returns('Org1MSP'),
        };
        ctx = {
            stub,
            clientIdentity: identity
        };
    });""",
    (12, 7): """    it('crea propiedad cuando no existe', async () => {
        // getState devuelve null (propiedad no existe)
        stub.getState.resolves(null);
        // putState acepta cualquier valor
        stub.putState.resolves();

        // Ejecutar
        const contract = new SmartContract();
        await contract.CreateProperty(ctx,
            'PROP001', 'Calle Mayor 1', 'Org1MSP',
            120, 'apartment', 250000);

        // Verificar
        expect(stub.putState.calledOnce).to.be.true;
    });
});""",
}


def replace_text_preserving_format(text_frame, new_text):
    first_run_format = None
    first_para_align = None
    for para in text_frame.paragraphs:
        if para.runs:
            first_run = para.runs[0]
            first_para_align = para.alignment
            rPr = first_run._r.find('{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
            if rPr is not None:
                first_run_format = deepcopy(rPr)
            break

    first_pPr = None
    for para in text_frame.paragraphs:
        pPr = para._pPr
        if pPr is not None:
            first_pPr = deepcopy(pPr)
            break

    txBody = text_frame._txBody
    for p in txBody.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}p'):
        txBody.remove(p)

    lines = new_text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
        else:
            p = text_frame.add_paragraph()

        if first_para_align is not None:
            p.alignment = first_para_align

        if first_pPr is not None:
            existing_pPr = p._pPr
            if existing_pPr is not None:
                p._p.remove(existing_pPr)
            p._p.insert(0, deepcopy(first_pPr))

        run = p.add_run()
        run.text = line
        if first_run_format is not None:
            existing_rPr = run._r.find('{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
            if existing_rPr is not None:
                run._r.remove(existing_rPr)
            run._r.insert(0, deepcopy(first_run_format))


def main():
    prs = Presentation(PPTX_PATH)
    for (slide_idx, shape_idx), new_text in REPLACEMENTS.items():
        try:
            slide = prs.slides[slide_idx]
            shape = slide.shapes[shape_idx]
            if not shape.has_text_frame:
                print(f"  ! Slide {slide_idx+1} shape {shape_idx}: no text_frame")
                continue
            replace_text_preserving_format(shape.text_frame, new_text)
            preview = new_text.split('\n')[0][:50]
            print(f"  OK Slide {slide_idx+1} shape {shape_idx}: '{preview}...'")
        except IndexError:
            print(f"  ! Slide {slide_idx+1} shape {shape_idx}: no existe")
    prs.save(PPTX_PATH)
    print(f"\nGuardado: {PPTX_PATH}")


if __name__ == "__main__":
    main()
