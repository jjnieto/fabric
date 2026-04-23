from pptx import Presentation

prs = Presentation("/mnt/d/Dev/Fabric/docs/slides/Modulo 4/dia_6.pptx")
targets = [8, 13]

for idx, slide in enumerate(prs.slides):
    if (idx+1) not in targets:
        continue
    print("=" * 60)
    print(f"SLIDE {idx+1}")
    print("=" * 60)
    for s_idx, shape in enumerate(slide.shapes):
        if shape.has_text_frame:
            paras = [p.text for p in shape.text_frame.paragraphs]
            combined = "\n".join(paras)
            try:
                l = shape.left/914400
                t = shape.top/914400
                w = shape.width/914400
                h = shape.height/914400
            except:
                l=t=w=h=0
            if combined.strip():
                first = combined.strip().split("\n")[0][:60]
                print(f"  Shape {s_idx} ({w:.1f}x{h:.1f}) pos=({l:.1f},{t:.1f}): '{first}'")
                if len(combined) > 100:
                    print(f"    FULL: {combined[:400]}")
    print()
