"""
Convierte codigo Go a Node.js en las slides del dia 4 del Modulo 4.
Abre el pptx existente y cambia SOLO el texto de los shapes identificados,
preservando formato (fuente, tamano, color, posicion, tamano del shape).
"""
from pptx import Presentation
from copy import deepcopy

PPTX_PATH = "/mnt/d/Dev/Fabric/docs/slides/Modulo 4/dia_4.pptx"

REPLACEMENTS = {
    # ============================================================
    # Slide 6 (idx 5): Token fungible: Mint
    # ============================================================
    (5, 4): "NODE.JS",
    (5, 6): """async Mint(ctx, amount) {
    // Solo el minter puede crear tokens
    const ok = ctx.clientIdentity
        .assertAttributeValue('role', 'minter');
    if (!ok) {
        throw new Error(
            'solo el minter puede emitir tokens');
    }
    if (parseInt(amount) <= 0) {
        throw new Error('amount debe ser positivo');
    }

    // Obtener ID del minter
    const minterID = ctx.clientIdentity.getID();""",
    (5, 7): """    // Actualizar balance
    const currentBalance = await getBalance(ctx, minterID);
    const newBalance = currentBalance + parseInt(amount);
    await ctx.stub.putState(balanceKey(minterID),
        Buffer.from(newBalance.toString()));

    // Actualizar totalSupply
    const supply = await getTotalSupply(ctx);
    await ctx.stub.putState('totalSupply',
        Buffer.from((supply + parseInt(amount)).toString()));

    // Evento
    ctx.stub.setEvent('Mint', Buffer.from(JSON.stringify(
        { minter: minterID, amount: parseInt(amount) })));
}""",

    # ============================================================
    # Slide 7 (idx 6): Token fungible: Transfer
    # ============================================================
    (6, 4): "NODE.JS",
    (6, 6): """async Transfer(ctx, to, amount) {
    amount = parseInt(amount);
    if (amount <= 0) {
        throw new Error('amount debe ser positivo');
    }

    const fromID = ctx.clientIdentity.getID();
    if (fromID === to) {
        throw new Error('no puedes transferirte a ti mismo');
    }

    // Verificar saldo
    const fromBalance = await getBalance(ctx, fromID);
    if (fromBalance < amount) {
        throw new Error(
            `saldo insuficiente: tiene ${fromBalance}, ` +
            `necesita ${amount}`);
    }""",
    (6, 7): """    // Actualizar balances
    const toBalance = await getBalance(ctx, to);
    await ctx.stub.putState(balanceKey(fromID),
        Buffer.from((fromBalance - amount).toString()));
    await ctx.stub.putState(balanceKey(to),
        Buffer.from((toBalance + amount).toString()));

    // Evento
    ctx.stub.setEvent('Transfer', Buffer.from(JSON.stringify({
        from: fromID, to, amount
    })));
}""",

    # ============================================================
    # Slide 11 (idx 10): NFT Mint y Transfer
    # ============================================================
    (10, 3): "NODE.JS",
    (10, 5): """// NFT en JSON (no hay struct en JS)
/**
 * @typedef {Object} NFT
 * @property {string} docType
 * @property {string} tokenID
 * @property {string} owner
 * @property {Object.<string,string>} metadata
 * @property {string} createdAt
 */

async MintNFT(ctx, tokenID, metadataJSON) {
    // Verificar que no existe
    const existing = await ctx.stub.getState('nft_' + tokenID);
    if (existing && existing.length > 0) {
        throw new Error(`NFT ${tokenID} ya existe`);
    }""",
    (10, 6): """    const ownerID = ctx.clientIdentity.getID();
    const metadata = JSON.parse(metadataJSON);

    const nft = {
        docType: 'nft',
        tokenID,
        owner: ownerID,
        metadata
    };

    // Guardar por tokenID y crear indice por owner
    await ctx.stub.putState('nft_' + tokenID,
        Buffer.from(JSON.stringify(nft)));
    const ownerKey = ctx.stub.createCompositeKey(
        'nftOwner', [ownerID, tokenID]);
    await ctx.stub.putState(ownerKey, Buffer.from([0x00]));
}""",

    # ============================================================
    # Slide 16 (idx 15): FoodLot estructura y creacion
    # ============================================================
    (15, 3): "NODE.JS",
    (15, 5): """// FoodLot: objeto JSON (no hay struct en JS)
/**
 * @typedef {Object} FoodLot
 * @property {string} docType
 * @property {string} lotID
 * @property {string} productType
 * @property {string} origin
 * @property {string} currentHolder
 * @property {string} status
 * @property {number} temperature
 * @property {number} weight
 * @property {HistEntry[]} history
 */

/**
 * @typedef {Object} HistEntry
 * @property {string} org
 * @property {string} action
 * @property {string} timestamp
 * @property {string} location
 */""",
    (15, 6): """async ProduceLot(ctx, lotID, productType, origin, weight) {
    const producerMSP = ctx.clientIdentity.getMSPID();

    const lot = {
        docType: 'foodlot',
        lotID,
        productType,
        origin,
        currentHolder: producerMSP,
        status: 'produced',
        weight: parseFloat(weight),
        history: [{
            org: producerMSP,
            action: 'produced',
            timestamp: new Date().toISOString()
        }]
    };

    await ctx.stub.putState('lot_' + lotID,
        Buffer.from(JSON.stringify(lot)));
}""",

    # ============================================================
    # Slide 17 (idx 16): FoodLot transferencia entre actores
    # ============================================================
    (16, 4): "NODE.JS",
    (16, 6): """async TransferLot(ctx, lotID, newHolder, location, temp) {
    const lot = await this.ReadLot(ctx, lotID);

    // Solo el holder actual puede transferir
    const callerMSP = ctx.clientIdentity.getMSPID();
    if (lot.currentHolder !== callerMSP) {
        throw new Error(
            `solo ${lot.currentHolder} puede transferir este lote`);
    }

    // Validar cold chain: alerta si temperatura fuera de rango
    const tempF = parseFloat(temp);
    if (tempF > 8.0) {
        lot.status = 'temp_alert';
        // No se bloquea, pero se registra la alerta
    } else {
        lot.status = 'in_transit';
    }""",
    (16, 7): """    // Actualizar holder y registrar en historial
    lot.currentHolder = newHolder;
    lot.temperature = tempF;
    lot.history.push({
        org: callerMSP,
        action: 'transferred_to_' + newHolder,
        timestamp: new Date().toISOString(),
        location
    });

    ctx.stub.setEvent('LotTransferred',
        Buffer.from(JSON.stringify({ lotID, to: newHolder })));

    await ctx.stub.putState('lot_' + lotID,
        Buffer.from(JSON.stringify(lot)));
}""",

    # ============================================================
    # Slide 18 (idx 17): FoodLot trazabilidad completa
    # ============================================================
    (17, 4): "NODE.JS",
    (17, 6): """// Consultar la trazabilidad completa de un lote
async GetLotHistory(ctx, lotID) {
    const lot = await this.ReadLot(ctx, lotID);
    return lot.history;
}

// Recall: retirar un lote del mercado
async RecallLot(ctx, lotID, reason) {
    // Solo el regulador puede hacer recall
    const ok = ctx.clientIdentity
        .assertAttributeValue('role', 'regulador');
    if (!ok) {
        throw new Error(
            'solo el regulador puede retirar lotes');
    }""",
    (17, 7): """    const lot = await this.ReadLot(ctx, lotID);

    lot.status = 'recalled';
    lot.history.push({
        org: 'Regulador',
        action: 'recalled: ' + reason,
        timestamp: new Date().toISOString()
    });

    ctx.stub.setEvent('LotRecalled',
        Buffer.from(JSON.stringify({ lotID, reason })));

    await ctx.stub.putState('lot_' + lotID,
        Buffer.from(JSON.stringify(lot)));
}""",
}


def replace_text_preserving_format(text_frame, new_text):
    """Reemplaza texto preservando formato del primer run."""
    first_run_format = None
    first_para_align = None
    for para in text_frame.paragraphs:
        if para.runs:
            first_run = para.runs[0]
            first_para_align = para.alignment
            rPr = first_run._r.find('{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
            if rPr is not None:
                first_run_format = deepcopy(rPr)
            break

    first_pPr = None
    for para in text_frame.paragraphs:
        pPr = para._pPr
        if pPr is not None:
            first_pPr = deepcopy(pPr)
            break

    txBody = text_frame._txBody
    for p in txBody.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}p'):
        txBody.remove(p)

    lines = new_text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
        else:
            p = text_frame.add_paragraph()

        if first_para_align is not None:
            p.alignment = first_para_align

        if first_pPr is not None:
            existing_pPr = p._pPr
            if existing_pPr is not None:
                p._p.remove(existing_pPr)
            p._p.insert(0, deepcopy(first_pPr))

        run = p.add_run()
        run.text = line
        if first_run_format is not None:
            existing_rPr = run._r.find('{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
            if existing_rPr is not None:
                run._r.remove(existing_rPr)
            run._r.insert(0, deepcopy(first_run_format))


def main():
    prs = Presentation(PPTX_PATH)

    for (slide_idx, shape_idx), new_text in REPLACEMENTS.items():
        try:
            slide = prs.slides[slide_idx]
            shape = slide.shapes[shape_idx]
            if not shape.has_text_frame:
                print(f"  ! Slide {slide_idx+1} shape {shape_idx}: no tiene text_frame, saltando")
                continue
            replace_text_preserving_format(shape.text_frame, new_text)
            preview = new_text.split('\n')[0][:50]
            print(f"  OK Slide {slide_idx+1} shape {shape_idx}: '{preview}...'")
        except IndexError:
            print(f"  ! Slide {slide_idx+1} shape {shape_idx}: no existe")

    prs.save(PPTX_PATH)
    print(f"\nGuardado: {PPTX_PATH}")


if __name__ == "__main__":
    main()
