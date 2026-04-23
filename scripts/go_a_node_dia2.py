"""
Convierte el codigo Go a Node.js en las slides del dia 2 del Modulo 4.
Abre el pptx existente y cambia SOLO el texto de los shapes identificados,
preservando el formato (fuente, tamano, color, posicion, tamano del shape).
"""
from pptx import Presentation
from copy import deepcopy
from lxml import etree

PPTX_PATH = "/mnt/d/Dev/Fabric/docs/slides/Modulo 4/dia_2.pptx"

# Mapping: (slide_index, shape_index) -> texto nuevo
# slide_index y shape_index son 0-based
REPLACEMENTS = {
    # ============================================================
    # Slide 6 (idx 5): Anatomia de un chaincode en Go -> Node.js
    # ============================================================
    (5, 1): "Anatomia de un chaincode en Node.js",  # titulo
    (5, 2): "Estructura basica: clase que extiende Contract",  # subtitulo
    (5, 4): "NODE.JS",  # badge
    (5, 6): """'use strict';
const { Contract } = require('fabric-contract-api');

// SmartContract proporciona las funciones del chaincode
class SmartContract extends Contract {

    // InitLedger inicializa el ledger con datos de ejemplo
    async InitLedger(ctx) {
        // Logica de inicializacion aqui
    }
}""",
    (5, 7): """// Exportar para que Fabric cargue el chaincode
module.exports = SmartContract;

// package.json:
// {
//   "main": "index.js",
//   "dependencies": {
//     "fabric-contract-api": "^2.2"
//   }
// }""",

    # ============================================================
    # Slide 11 (idx 10): Composite Keys: ejemplo practico
    # ============================================================
    (10, 3): "NODE.JS",  # badge (shape 3 en esta slide)
    (10, 5): """// Crear una clave compuesta para un activo
async CreateAsset(ctx, owner, assetType, id, value) {
    // Clave compuesta: permite buscar por owner o por type
    const key = ctx.stub.createCompositeKey('asset',
        [owner, assetType, id]);

    const asset = { ID: id, Owner: owner,
                    Type: assetType, Value: value };

    await ctx.stub.putState(key,
        Buffer.from(JSON.stringify(asset)));
}""",
    (10, 6): """// Buscar todos los activos de un propietario
async GetAssetsByOwner(ctx, owner) {
    const iterator = await ctx.stub
        .getStateByPartialCompositeKey('asset', [owner]);
    // ... iterar sobre resultados
}""",

    # ============================================================
    # Slide 13 (idx 12): Modelo de datos Registro de Propiedad
    # ============================================================
    (12, 4): "NODE.JS",  # badge
    (12, 6): """// Propiedad inmobiliaria (JSDoc describe los campos)
/**
 * @typedef {Object} Property
 * @property {string} docType        - siempre "property"
 * @property {string} id             - referencia catastral
 * @property {string} address        - direccion fisica
 * @property {string} owner          - ID del propietario (MSP)
 * @property {string} ownerName      - nombre del propietario
 * @property {number} area           - metros cuadrados
 * @property {string} propertyType   - "apartment"|"house"|"land"
 * @property {string} status         - "active"|"transferring"|"blocked"
 * @property {string} registeredAt   - fecha de registro
 * @property {number} appraisalValue - valor de tasacion (euros)
 */

// Key design: "property" + comunidad + municipio + referencia_catastral
// Permite buscar por comunidad, por municipio, o por referencia exacta""",

    # ============================================================
    # Slide 15 (idx 14): Rich Queries con CouchDB
    # ============================================================
    (14, 4): "NODE.JS",  # badge
    (14, 6): """// Buscar propiedades por tipo y rango de valor
async QueryProperties(ctx, propertyType, minValue) {
    // Mango query (sintaxis CouchDB)
    const queryString = JSON.stringify({
        selector: {
            docType: 'property',
            propertyType: propertyType,
            appraisalValue: { $gte: parseInt(minValue) }
        },
        sort: [{ appraisalValue: 'desc' }],
        use_index: 'indexPropertyType'
    });

    const iterator = await ctx.stub.getQueryResult(queryString);
    // Iterar y deserializar resultados...
}""",

    # ============================================================
    # Slide 17 (idx 16): Patron CRUD: Create
    # ============================================================
    (16, 4): "NODE.JS",  # badge
    (16, 6): """async CreateProperty(ctx, id, address, owner,
    area, propType, value) {

    // Verificar que no existe
    const existing = await ctx.stub.getState(id);
    if (existing && existing.length > 0) {
        throw new Error(
            `property ${id} already exists`);
    }""",
    (16, 7): """    const property = {
        docType: 'property',
        id,
        address,
        owner,
        area: parseInt(area),
        propertyType: propType,
        status: 'active',
        appraisalValue: parseInt(value)
    };

    await ctx.stub.putState(id,
        Buffer.from(JSON.stringify(property)));
}""",

    # ============================================================
    # Slide 18 (idx 17): Patron CRUD: Read
    # ============================================================
    (17, 4): "NODE.JS",  # badge
    (17, 6): """async ReadProperty(ctx, id) {
    const propertyJSON = await ctx.stub.getState(id);

    if (!propertyJSON || propertyJSON.length === 0) {
        throw new Error(
            `property ${id} does not exist`);
    }

    const property = JSON.parse(propertyJSON.toString());
    return property;
}""",
    (17, 7): """// En Go seria:
// func (s *SmartContract) ReadProperty(
//     ctx ..., id string) (*Property, error) {
//     data, err := ctx.GetStub().GetState(id)
//     if err != nil { return nil, err }
//     if data == nil {
//         return nil, errors.New("not found")
//     }
//     var p Property
//     json.Unmarshal(data, &p)
//     return &p, nil
// }""",

    # ============================================================
    # Slide 19 (idx 18): Patron CRUD: Update
    # ============================================================
    (18, 4): "NODE.JS",  # badge
    (18, 6): """async UpdatePropertyValue(ctx, id, newValue) {
    // Leer el activo existente
    const property = await this.ReadProperty(ctx, id);

    // Verificar permisos: solo el owner puede actualizar
    const clientMSP = ctx.clientIdentity.getMSPID();
    if (property.owner !== clientMSP) {
        throw new Error(
            'only the owner can update this property');
    }""",
    (18, 7): """    // Actualizar campo
    property.appraisalValue = parseInt(newValue);

    // Guardar
    await ctx.stub.putState(id,
        Buffer.from(JSON.stringify(property)));
}""",

    # ============================================================
    # Slide 20 (idx 19): Patron CRUD: Delete y GetAll
    # ============================================================
    (19, 4): "NODE.JS",  # badge
    (19, 6): """// Delete: borrado logico (cambiar status)
//         vs fisico (deleteState)
async DeactivateProperty(ctx, id) {
    const property = await this.ReadProperty(ctx, id);

    // borrado logico: mantiene historial
    property.status = 'inactive';

    await ctx.stub.putState(id,
        Buffer.from(JSON.stringify(property)));
}""",
    (19, 7): """// GetAll con paginacion
async GetAllProperties(ctx, pageSize, bookmark) {
    const { iterator, metadata } = await ctx.stub
        .getStateByRangeWithPagination('', '',
            parseInt(pageSize), bookmark);
    // ... iterar y construir resultado
    //     con metadata.bookmark
}""",
}


def replace_text_preserving_format(text_frame, new_text):
    """
    Reemplaza el texto del text_frame preservando el formato del primer run.
    Cada linea de new_text se convierte en un paragraph.
    """
    # Capturar formato del primer run (si existe)
    first_run_format = None
    first_para_align = None
    for para in text_frame.paragraphs:
        if para.runs:
            first_run = para.runs[0]
            first_para_align = para.alignment
            # Extraer el elemento rPr (run properties) como XML para copiarlo
            rPr = first_run._r.find('{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
            if rPr is not None:
                first_run_format = deepcopy(rPr)
            break

    # Capturar formato del parrafo (pPr) si existe
    first_pPr = None
    for para in text_frame.paragraphs:
        pPr = para._pPr
        if pPr is not None:
            first_pPr = deepcopy(pPr)
            break

    # Limpiar todos los parrafos
    txBody = text_frame._txBody
    for p in txBody.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}p'):
        txBody.remove(p)

    # Anadir nuevos parrafos, uno por linea
    lines = new_text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
        else:
            p = text_frame.add_paragraph()

        # Aplicar alignment del parrafo original si existe
        if first_para_align is not None:
            p.alignment = first_para_align

        # Aplicar pPr del parrafo original (copia)
        if first_pPr is not None:
            existing_pPr = p._pPr
            if existing_pPr is not None:
                p._p.remove(existing_pPr)
            p._p.insert(0, deepcopy(first_pPr))

        # Anadir un run con el texto y el formato
        run = p.add_run()
        run.text = line
        if first_run_format is not None:
            # Copiar el rPr del primer run original
            existing_rPr = run._r.find('{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
            if existing_rPr is not None:
                run._r.remove(existing_rPr)
            run._r.insert(0, deepcopy(first_run_format))


def main():
    prs = Presentation(PPTX_PATH)

    for (slide_idx, shape_idx), new_text in REPLACEMENTS.items():
        try:
            slide = prs.slides[slide_idx]
            shape = slide.shapes[shape_idx]
            if not shape.has_text_frame:
                print(f"  ! Slide {slide_idx+1} shape {shape_idx}: no tiene text_frame, saltando")
                continue
            replace_text_preserving_format(shape.text_frame, new_text)
            preview = new_text.split('\n')[0][:50]
            print(f"  OK Slide {slide_idx+1} shape {shape_idx}: '{preview}...'")
        except IndexError:
            print(f"  ! Slide {slide_idx+1} shape {shape_idx}: no existe, saltando")

    prs.save(PPTX_PATH)
    print(f"\nGuardado: {PPTX_PATH}")


if __name__ == "__main__":
    main()
