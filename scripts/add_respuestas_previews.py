"""
Anade slides de respuestas a los 5 previews del Modulo 3.
Cada preview recibe 2 slides: respuestas a preguntas clave + respuestas al debate.
"""
import sys
sys.path.insert(0, "/mnt/d/Dev/Fabric/scripts")
from gen_helpers import add_content_slide
from pptx import Presentation

BASE = "/mnt/d/Dev/Fabric/docs/slides/Modulo 3"

# ============================================================
# PREVIEW 01 - WALMART
# ============================================================
prs = Presentation(f"{BASE}/preview01_walmart.pptx")

add_content_slide(prs, "Respuestas a las preguntas clave", [
    "1. ¿Cuantas organizaciones y con peer propio?",
    "2. ¿Un canal unico o multiples canales segregados?",
    "3. ¿Datos on-chain vs Private Data?",
    "4. ¿Quien puede crear un lote?",
    "5. ¿Quien puede hacer recall?",
    "6. Politica de endorsement (normal vs recall):",
    "7. ¿Como verifica el consumidor final un QR?",
], subbullets={
    0: ["4 orgs (Productor, Distribuidor, Supermercado, Regulador), cada una con su peer."],
    1: ["Un canal unico compartido — la trazabilidad NECESITA ser publica entre miembros."],
    2: ["On-chain: lotID, origen, estado, movimientos (historial).",
        "Private Data: precios entre actores, margenes, terminos comerciales."],
    3: ["Solo el Productor (validar MSPID == ProductorMSP en el chaincode)."],
    4: ["Solo el Regulador (via atributo role=regulador en cert X.509 — ABAC)."],
    5: ["Movimientos normales: solo el holder actual endorsa (state-based endorsement).",
        "Recall: AND(Productor, Distribuidor, Supermercado) para consenso amplio."],
    6: ["Web publica del consorcio con QR -> API REST -> chaincode query.",
        "El consumidor no tiene peer ni certificado — es un lector publico."],
})

add_content_slide(prs, "Respuestas al debate", [
    "1. ¿Vuestra topologia vs IBM Food Trust?",
    "2. ¿Ventaja de Fabric vs una BD gestionada por Walmart?",
    "3. ¿Y si el productor miente sobre el origen?",
    "4. Walmart impuso la adopcion. ¿Justo? ¿Necesario?",
    "5. Integrar sensores IoT para temperatura:",
], subbullets={
    0: ["Food Trust tiene +150 orgs; nuestra red es reducida para aula.",
        "Diferencia clave: Food Trust no tiene regulador como org — Walmart es el fundador dominante."],
    1: ["Walmart gestionando la BD: los competidores (Carrefour, Lidl) desconfian.",
        "Fabric: nadie controla solo; cada org valida independientemente."],
    2: ["Blockchain garantiza INTEGRIDAD, no VERACIDAD.",
        "Si el productor miente al registrar, la mentira queda inmutable.",
        "Mitigar: auditorias fisicas, sensores IoT, validacion cruzada, responsabilidad legal."],
    3: ["Funcional (es el lider), pero no ideal — genera problema del fundador dominante.",
        "Ideal: un regulador (AESAN, FDA) impulsa la adopcion; el mercado se suma.",
        "Aprender de TradeLens: sin los competidores no hay masa critica."],
    4: ["Gateway IoT intermedio con certificado X.509.",
        "Los sensores firman con HMAC/TPM; el gateway valida y envia al chaincode.",
        "El chaincode ve: 'gateway-X certifica que sensor-Y reporto temp=X°C a las H'."],
})

prs.save(f"{BASE}/preview01_walmart.pptx")
print(f"preview01_walmart.pptx: total {len(prs.slides)} slides")

# ============================================================
# PREVIEW 02 - WETRADE
# ============================================================
prs = Presentation(f"{BASE}/preview02_wetrade.pptx")

add_content_slide(prs, "Respuestas a las preguntas clave", [
    "1. ¿Las PYMEs: orgs o usuarios?",
    "2. ¿Un canal o canales por pareja de bancos?",
    "3. ¿Importes y datos sensibles?",
    "4. Estados de una operacion:",
    "5. Politicas de endorsement (crear vs liberar pago):",
    "6. ¿BCE en el canal?",
    "7. ¿Escala N*(N-1)/2 Private Data Collections?",
], subbullets={
    0: ["Usuarios de los bancos (cert X.509 con role=client emitido por el banco).",
        "Hacer orgs a cada PYME seria inviable: miles/millones de orgs."],
    1: ["Un canal compartido con Private Data Collections por pareja de bancos.",
        "Permite detectar fraudes transversales (doble financiacion, lavado)."],
    2: ["Private Data Collections por pareja de bancos involucrados.",
        "Solo los 2 bancos ven el importe; los demas solo ven el hash."],
    3: ["proposed -> approved -> in_transit -> delivered -> paid",
        "Con transicion rechazada / disputa en cualquier punto."],
    4: ["Crear operacion: AND(BancoVendedor, BancoComprador).",
        "Liberar pago: AND(BancoVendedor, BancoComprador) + verificar status='delivered'."],
    5: ["Si como org de auditoria (solo lectura), NO como endorser.",
        "Evita que el regulador pueda bloquear operaciones comerciales."],
    6: ["No escala para N grandes: con 50 bancos son 1225 PDCs.",
        "Alternativa: state-based endorsement dinamico por operacion.",
        "We.Trade llego al limite con 12 bancos — esto contribuyo al cierre."],
})

add_content_slide(prs, "Respuestas al debate", [
    "1. ¿Por que cerro We.Trade a pesar de tener 12 bancos?",
    "2. ¿Que modelo de ingresos propondriais?",
    "3. ¿PYMEs como usuarios tiene sentido?",
    "4. PYME cambia de banco: ¿y su historial?",
    "5. ¿Regulador (BCE) en el canal como HKMA?",
], subbullets={
    0: ["Problema de modelo de negocio: nadie queria financiar indefinidamente.",
        "Los miembros querian participar pero no pagar cuotas sostenidas.",
        "Faltaba un flujo de ingresos recurrente ligado al uso real."],
    1: ["Fees por transaccion: 0.1-0.5% del importe financiado.",
        "Cuota anual escalada: pequenos bancos pagan menos que grandes.",
        "Modelo hibrido: fee por uso + cuota minima para cubrir infra."],
    2: ["Si — las PYMEs son clientes de los bancos, no actores de la red.",
        "Ventaja: escala infinita sin anadir orgs.",
        "Cambia: la identidad de la PYME la emite su banco, que es responsable."],
    3: ["Su historial queda en el ledger (asociado a cert viejo).",
        "El nuevo banco emite un cert nuevo con el mismo ID (p.ej. CIF).",
        "Funcion admin: 'link_identity(old_cert_hash, new_cert)' — auditada."],
    4: ["Ventajas: supervision en tiempo real, deteccion temprana de fraude.",
        "Riesgos: concentracion de informacion, GDPR, privacidad bancaria.",
        "Compromise: BCE como org con acceso de SOLO LECTURA en un canal auxiliar."],
})

prs.save(f"{BASE}/preview02_wetrade.pptx")
print(f"preview02_wetrade.pptx: total {len(prs.slides)} slides")

# ============================================================
# PREVIEW 03 - HKMA
# ============================================================
prs = Presentation(f"{BASE}/preview03_hkma.pptx")

add_content_slide(prs, "Respuestas a las preguntas clave", [
    "1. ¿Regulador endorsa o solo audita?",
    "2. Identificar univocamente una factura:",
    "3. Datos personales on-chain (GDPR):",
    "4. ¿Regulador puede CONGELAR operaciones?",
    "5. Si el cert del regulador es comprometido:",
    "6. Rich query por invoiceHash: ¿LevelDB o CouchDB?",
], subbullets={
    0: ["Solo audita. El regulador NO endorsea transacciones normales.",
        "Evita que el regulador pueda alterar el flujo de negocio.",
        "Excepcion: funciones administrativas (bloquear, marcar) si endorsa."],
    1: ["Hash SHA-256 del contenido de la factura (datos normalizados).",
        "El hash es unico, determinista y no revela el contenido."],
    2: ["NO guardar datos personales (nombres, DNIs, direcciones) on-chain.",
        "Guardar solo HASHES o IDs pseudonimos.",
        "Los datos reales van off-chain (BD tradicional de cada banco)."],
    3: ["Si, pero con garantias:",
        "  - Evento de auditoria registrado (quien, cuando, por que)",
        "  - Los bancos involucrados tienen derecho a apelar",
        "  - Bloquear != confiscar; el estado del activo cambia, no desaparece."],
    4: ["Revocar el cert via Fabric CA + generar CRL + distribuir a peers.",
        "Hasta que los peers aplican la CRL, el cert comprometido es valido.",
        "Mitigacion: rotacion periodica + monitoreo de comportamiento anomalo."],
    5: ["LevelDB: rich queries NO funcionan, habria que escanear todo.",
        "CouchDB: indice por invoiceHash permite busqueda O(log n).",
        "Respuesta: CouchDB es IMPRESCINDIBLE para este caso de uso."],
})

add_content_slide(prs, "Respuestas al debate", [
    "1. ¿Regulador en el consorcio o acceso via API?",
    "2. En Europa, ¿BCE o banco central nacional?",
    "3. ¿Regulador con acceso total es privacy-friendly?",
    "4. Si el regulador es comprometido:",
    "5. ¿Puede un banco oponerse a una marca del regulador?",
], subbullets={
    0: ["En el consorcio = supervision en tiempo real, deteccion inmediata.",
        "Via API = mas control de acceso, menos intrusivo, mas lento.",
        "Lo correcto depende del nivel de confianza y del marco legal local."],
    1: ["BCE a nivel europeo — operaciones transfronterizas.",
        "Banco central nacional para operaciones domesticas (p.ej. Banco de Espana).",
        "Modelo hibrido: federacion de bancos centrales (similar a EBSI)."],
    2: ["No es privacy-friendly por naturaleza, pero SI es 'auditoria razonable'.",
        "El regulador ya tiene acceso legal a esta informacion (aunque mas lento).",
        "Clave: que sea PROPORCIONAL y AUDITABLE (acceso justificado y registrado)."],
    3: ["Grave — el regulador tiene visibilidad total de operaciones.",
        "Requiere: rotacion frecuente de cert, HSM para la clave privada,",
        "multi-factor para operaciones sensibles, auditoria de los accesos del regulador."],
    4: ["Legalmente: via tribunales (el regulador tiene autoridad limitada).",
        "Tecnicamente: NO dentro del chaincode (la politica lo rechaza).",
        "El banco puede solicitar revision, pero la operacion queda marcada hasta aclarar."],
})

prs.save(f"{BASE}/preview03_hkma.pptx")
print(f"preview03_hkma.pptx: total {len(prs.slides)} slides")

# ============================================================
# PREVIEW 04 - TRADELENS
# ============================================================
prs = Presentation(f"{BASE}/preview04_tradelens.pptx")

add_content_slide(prs, "Respuestas a las preguntas clave", [
    "1. ¿Quien es el fundador ideal?",
    "2. Orderers distribuidos:",
    "3. Modelo de financiacion:",
    "4. Incentivo para MSC, CMA CGM:",
    "5. ¿Puertos y aduanas: orgs plenas o limitadas?",
    "6. Transbordo A -> B: ¿quien endorsa?",
], subbullets={
    0: ["Fundacion neutral sin animo de lucro (modelo Alastria, Hyperledger).",
        "IMO (Organizacion Maritima Internacional) seria la opcion institucional.",
        "Lo que NO: una naviera o el propio IBM."],
    1: ["Raft con 3-5 nodos, cada uno en una naviera diferente.",
        "Rotacion anual del 'leader' para evitar concentracion.",
        "Ningun nodo puede ser apagado sin consentimiento de la fundacion."],
    2: ["Cuota escalada por TEU transportados (uso real).",
        "Cuota minima para cubrir infra (incluso a navieras pequenas).",
        "Incentivo: los primeros 3 anos, fees reducidos para los fundadores."],
    3: ["Visibilidad completa de trazabilidad (mas eficientes en transbordos).",
        "Menor papeleo (digitalizacion de Bill of Lading).",
        "Cumplimiento de regulaciones automatico (aduanas).",
        "Reduccion de fraudes (no pueden manipular datos ajenos)."],
    4: ["Orgs plenas — son parte operativa del flujo (no solo lectores).",
        "Puertos: escriben eventos (loading, discharge) y validan.",
        "Aduanas: escriben eventos (cleared) con rol especial de autoridad."],
    5: ["Politica AND(Naviera_A, Naviera_B) — ambas deben firmar el transbordo.",
        "Se implementa con state-based endorsement por contenedor.",
        "Si una naviera no firma, el transbordo queda pendiente (no se commitea)."],
})

add_content_slide(prs, "Respuestas al debate", [
    "1. ¿Por que fundacion > 'Maersk lanza la plataforma'?",
    "2. ¿Aceptaria Maersk participar con 1 voto = 1 miembro?",
    "3. ¿Es realista que competidores directos colaboren?",
    "4. Rol de la IMO en tu solucion:",
    "5. ¿Quien deberia relanzar este proyecto?",
], subbullets={
    0: ["La fundacion es neutral y no tiene incentivos comerciales propios.",
        "Maersk, por definicion, es competidor de los demas — genera desconfianza.",
        "La historia demuestra: cuando el fundador es competidor, los demas no se suman."],
    1: ["Probablemente no sin incentivos. Necesita: beneficios operativos claros,",
        "reputacion por liderar sostenibilidad, reduccion de costes auditables."],
    2: ["Si — ya lo hacen en IATA, GSM Association, Visa (antes consorcio).",
        "La clave es alinear incentivos: mas valor conjunto que el individual."],
    3: ["Organismo rector, establece estandares tecnicos y politicos.",
        "Puede dar legitimidad internacional y apoyo de gobiernos.",
        "No opera la red (delegado a la fundacion), pero la supervisa."],
    4: ["Una coalicion internacional de navieras + IMO + grandes puertos.",
        "Lanzar con MVP pequeno (2-3 navieras neutrales) y crecer.",
        "Financiacion mixta: fondos de los miembros + apoyo institucional."],
})

prs.save(f"{BASE}/preview04_tradelens.pptx")
print(f"preview04_tradelens.pptx: total {len(prs.slides)} slides")

# ============================================================
# PREVIEW 05 - B3i
# ============================================================
prs = Presentation(f"{BASE}/preview05_b3i.pptx")

add_content_slide(prs, "Respuestas a las preguntas clave", [
    "1. Fee por transaccion, cuota anual o escalado:",
    "2. Contratos multi-parte: ¿TODAS o MAJORITY?",
    "3. Validar que un siniestro es real:",
    "4. EIOPA: ¿voto o solo lectura?",
    "5. Datos del cliente final on-chain (GDPR):",
    "6. Contratos viejos pre-blockchain:",
], subbullets={
    0: ["Mix: cuota anual (cubre infra) + fee por transaccion (valor por uso).",
        "Escalado por volumen: mas contratos = menor % por contrato.",
        "Clave: que el fee sea marginal vs el ahorro operativo (< 5%)."],
    1: ["TODAS las partes firmantes — cada reasegurador debe aprobar.",
        "State-based endorsement por contrato: la politica la establecen las partes.",
        "MAJORITY seria peligroso: permite liquidar sin el reasegurador principal."],
    2: ["Oraculos externos (datos meteo, sismicos, noticias).",
        "Peritaje por empresa certificada (rol auditor con cert X.509).",
        "Validacion cruzada: aseguradora primaria + reasegurador independiente.",
        "Combinacion de las tres: mas costosa pero mas robusta."],
    3: ["Solo lectura para auditoria + capacidad de bloquear operaciones sospechosas.",
        "NO voto en governance: mantiene independencia.",
        "Interviene solo cuando hay indicios de fraude o incumplimiento."],
    4: ["NO datos personales on-chain. Solo hashes o IDs pseudonimos.",
        "Los datos reales van off-chain en cada aseguradora.",
        "Cumplimiento del derecho al olvido: borrar off-chain hace el hash huerfano."],
    5: ["Migracion gradual — mantener sistema viejo para contratos existentes.",
        "Nuevos contratos desde X fecha en blockchain.",
        "Ventana de transicion: 2-3 anos para renovar contratos existentes."],
})

add_content_slide(prs, "Respuestas al debate", [
    "1. TradeLens (gobernanza) vs B3i (modelo de negocio): ¿cual es peor?",
    "2. ¿Cual es mas facil de arreglar a posteriori?",
    "3. ¿Cual evitarias PRIMERO en tu proyecto?",
    "4. ¿Es posible evitar ambos al mismo tiempo?",
    "5. Si relanzas B3i, ¿que harias distinto?",
], subbullets={
    0: ["Gobernanza (TradeLens) es mas grave — sin participantes, no hay red.",
        "Modelo de negocio (B3i) es recuperable — hay red operativa, falta solo financiacion.",
        "En ambos casos, la tecnologia funcionaba."],
    1: ["Modelo de negocio — se puede cambiar pricing sin rediseno tecnico.",
        "Gobernanza requiere recrear la red o reconstruir la base de miembros.",
        "Por eso B3i es recuperable; TradeLens no se pudo reflotar."],
    2: ["Gobernanza — si no es neutral desde el dia 1, nunca lo sera.",
        "Modelo de negocio — se puede iterar (MVP -> validar -> ajustar precio)."],
    3: ["Si — se hace, pero requiere diseño cuidadoso:",
        "  Gobernanza neutral (fundacion, rotacion, no fundador dominante)",
        "  Modelo de ingresos ligado al uso real (fees por transaccion)",
        "  MVP pequeno con casos claros (antes de pretender cubrirlo todo)."],
    4: ["Empezar con MVP muy pequeno: 2-3 aseguradoras, una linea (catastroficos).",
        "Modelo de ingresos claro ANTES de desarrollar (pilot fees por contrato).",
        "Gobernanza via fundacion sin animo de lucro.",
        "Integracion progresiva con sistemas legacy (no big-bang)."],
})

prs.save(f"{BASE}/preview05_b3i.pptx")
print(f"preview05_b3i.pptx: total {len(prs.slides)} slides")
