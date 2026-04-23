from pptx import Presentation

for fname in ["dia_2.pptx", "dia_4.pptx"]:
    path = f"/mnt/d/Dev/Fabric/docs/slides/Modulo 4/{fname}"
    prs = Presentation(path)
    print("=" * 70)
    print(fname)
    print("=" * 70)
    for idx, slide in enumerate(prs.slides):
        has_go_badge = False
        has_go_title = False
        has_go_code = False
        title = ""
        shapes_info = []
        for s_idx, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t and not title and len(t) > 5:
                        title = t[:80]
                full = "\n".join([p.text for p in shape.text_frame.paragraphs])
                if full.strip() == "GO":
                    has_go_badge = True
                    shapes_info.append((s_idx, "BADGE_GO"))
                if "func " in full and "ctx contract" in full:
                    has_go_code = True
                    shapes_info.append((s_idx, "GO_CODE"))
                if "type " in full and "struct" in full:
                    has_go_code = True
                    shapes_info.append((s_idx, "GO_STRUCT"))
        if "en Go" in title.lower() or "Anatomia" in title and "Go" in title:
            has_go_title = True

        if has_go_badge or has_go_code or has_go_title:
            print(f"\nSlide {idx+1}: {title}")
            for s_idx, kind in shapes_info:
                shape = slide.shapes[s_idx]
                first = ""
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        first = para.text.strip()[:60]
                        break
                print(f"  Shape {s_idx} [{kind}]: '{first}'")
    print()
