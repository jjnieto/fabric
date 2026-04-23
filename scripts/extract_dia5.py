from pptx import Presentation

prs = Presentation("/mnt/d/Dev/Fabric/docs/slides/Modulo 4/dia_5.pptx")
targets = [7, 9, 13, 17]

for idx, slide in enumerate(prs.slides):
    if (idx+1) not in targets:
        continue
    print("=" * 60)
    print(f"SLIDE {idx+1}")
    print("=" * 60)
    for s_idx, shape in enumerate(slide.shapes):
        if shape.has_text_frame:
            paras = [p.text for p in shape.text_frame.paragraphs]
            combined = "\n".join(paras)
            if len(combined.strip()) > 50:
                try:
                    w = shape.width/914400
                    h = shape.height/914400
                except:
                    w=h=0
                print(f"\n--- Shape {s_idx} ({w:.1f}x{h:.1f}) ---")
                print(combined)
    print()
