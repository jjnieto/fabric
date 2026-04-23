from pptx import Presentation

for fname in ["dia_5.pptx", "dia_6.pptx"]:
    path = f"/mnt/d/Dev/Fabric/docs/slides/Modulo 4/{fname}"
    prs = Presentation(path)
    print("=" * 70)
    print(fname)
    print("=" * 70)
    for idx, slide in enumerate(prs.slides):
        shapes_info = []
        title = ""
        for s_idx, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t and not title and len(t) > 5:
                        title = t[:80]
                full = "\n".join([p.text for p in shape.text_frame.paragraphs])
                if full.strip() == "GO":
                    shapes_info.append((s_idx, "BADGE_GO", ""))
                elif "func " in full and ("ctx contract" in full or "SmartContract" in full):
                    first = full.split("\n")[0].strip()[:70]
                    shapes_info.append((s_idx, "GO_CODE", first))
                elif "type " in full and "struct" in full:
                    first = full.split("\n")[0].strip()[:70]
                    shapes_info.append((s_idx, "GO_STRUCT", first))
        if shapes_info:
            print(f"\nSlide {idx+1}: {title}")
            for s_idx, kind, first in shapes_info:
                shape = slide.shapes[s_idx]
                try:
                    l = shape.left/914400
                    t = shape.top/914400
                    w = shape.width/914400
                    h = shape.height/914400
                except:
                    l=t=w=h=0
                print(f"  Shape {s_idx} [{kind}] ({w:.1f}x{h:.1f}): {first}")
    print()
