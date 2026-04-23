"""
Convierte Go a Node.js en slides del dia 5 del Modulo 4.
Abre pptx existente y modifica solo texto de shapes identificados.
"""
from pptx import Presentation
from copy import deepcopy

PPTX_PATH = "/mnt/d/Dev/Fabric/docs/slides/Modulo 4/dia_5.pptx"

REPLACEMENTS = {
    # Slide 7 (idx 6): State-Based Endorsement: ejemplo
    (6, 4): "NODE.JS",
    (6, 6): """const { KeyEndorsementPolicy } = require('fabric-shim');

async SetHighValueEndorsement(ctx, propertyID) {
    // Leer propiedad
    const property = await this.ReadProperty(ctx, propertyID);

    // Si el valor supera 1M, requerir endorsement de 3 orgs
    if (property.appraisalValue > 1000000) {
        const policy = new KeyEndorsementPolicy();""",
    (6, 7): """        // Requiere las tres organizaciones
        policy.addOrgs('MEMBER',
            'Org1MSP', 'Org2MSP', 'Org3MSP');
        const policyBytes = policy.getPolicy();

        await ctx.stub.setStateValidationParameter(
            propertyID, policyBytes);
    }
}""",

    # Slide 9 (idx 8): ABAC con multiples roles
    (8, 4): "NODE.JS",
    (8, 6): """async checkRole(ctx, requiredRole) {
    const [role, found] = ctx.clientIdentity
        .getAttributeValue('role');

    if (!found) {
        throw new Error(
            "el certificado no tiene atributo 'role'");
    }

    if (role !== requiredRole) {
        throw new Error(
            `rol requerido: ${requiredRole}, ` +
            `rol actual: ${role}`);
    }
}""",
    (8, 7): """// Uso en funciones del chaincode:
async ApproveTransfer(ctx, ...args) {
    await this.checkRole(ctx, 'registrador');
    // ... logica de aprobacion
}

async FreezeProperty(ctx, ...args) {
    await this.checkRole(ctx, 'autoridad_judicial');
    // ... logica de bloqueo
}""",

    # Slide 13 (idx 12): Determinismo ejemplos
    (12, 4): "NODE.JS",
    (12, 6): """// ERROR 1: Timestamp del sistema
async BadCreate(ctx) {
    asset.createdAt = new Date().toISOString();
    // DIFERENTE EN CADA PEER
    // Correcto: usar ctx.stub.getTxTimestamp()
}

// ERROR 2: Iterar un objeto (orden no garantizado)
async BadTotals(ctx) {
    const totals = { a: 1, b: 2, c: 3 };
    let result = '';
    for (const k in totals) {
        result += `${k}:${totals[k]},`;
    }
    // result puede ser "a:1,b:2,c:3" o cualquier otro orden
    // -> NO DETERMINISTA (depende del motor JS)
}""",
    (12, 7): """// ERROR 3: Llamada HTTP externa
async BadPrice(ctx) {
    const resp = await fetch(
        'https://api.prices.com/gold');
    // Cada peer puede obtener un precio diferente
    // -> INVALIDA
}

// ERROR 4: Numero aleatorio
async BadRandom(ctx) {
    const id = `ASSET-${Math.floor(Math.random()*99999)}`;
    // ID diferente en cada peer -> INVALIDA
}""",

    # Slide 17 (idx 16): Migracion de datos
    (16, 4): "NODE.JS",
    (16, 6): """// Version 1 del modelo (documentacion JSDoc)
/**
 * @typedef {Object} PropertyV1
 * @property {string} id
 * @property {string} owner
 * @property {string} address
 */

// Version 2: nuevos campos
/**
 * @typedef {Object} PropertyV2
 * @property {string} id
 * @property {string} owner
 * @property {string} address
 * @property {string} propertyType    - NUEVO
 * @property {number} appraisalValue  - NUEVO
 * @property {number} version         - NUEVO: control version
 */""",
    (16, 7): """// Leer con compatibilidad: si faltan campos, defaults
async ReadProperty(ctx, id) {
    const data = await ctx.stub.getState(id);
    const prop = JSON.parse(data.toString());

    // Migracion lazy: si es version antigua, aplicar defaults
    if ((prop.version || 0) < 2) {
        if (!prop.propertyType) {
            prop.propertyType = 'unknown';
        }
        prop.version = 2;
        // Guardar migrado (opcional: on-read o funcion aparte)
    }
    return prop;
}""",
}


def replace_text_preserving_format(text_frame, new_text):
    first_run_format = None
    first_para_align = None
    for para in text_frame.paragraphs:
        if para.runs:
            first_run = para.runs[0]
            first_para_align = para.alignment
            rPr = first_run._r.find('{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
            if rPr is not None:
                first_run_format = deepcopy(rPr)
            break

    first_pPr = None
    for para in text_frame.paragraphs:
        pPr = para._pPr
        if pPr is not None:
            first_pPr = deepcopy(pPr)
            break

    txBody = text_frame._txBody
    for p in txBody.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}p'):
        txBody.remove(p)

    lines = new_text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
        else:
            p = text_frame.add_paragraph()

        if first_para_align is not None:
            p.alignment = first_para_align

        if first_pPr is not None:
            existing_pPr = p._pPr
            if existing_pPr is not None:
                p._p.remove(existing_pPr)
            p._p.insert(0, deepcopy(first_pPr))

        run = p.add_run()
        run.text = line
        if first_run_format is not None:
            existing_rPr = run._r.find('{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
            if existing_rPr is not None:
                run._r.remove(existing_rPr)
            run._r.insert(0, deepcopy(first_run_format))


def main():
    prs = Presentation(PPTX_PATH)
    for (slide_idx, shape_idx), new_text in REPLACEMENTS.items():
        try:
            slide = prs.slides[slide_idx]
            shape = slide.shapes[shape_idx]
            if not shape.has_text_frame:
                print(f"  ! Slide {slide_idx+1} shape {shape_idx}: no text_frame")
                continue
            replace_text_preserving_format(shape.text_frame, new_text)
            preview = new_text.split('\n')[0][:50]
            print(f"  OK Slide {slide_idx+1} shape {shape_idx}: '{preview}...'")
        except IndexError:
            print(f"  ! Slide {slide_idx+1} shape {shape_idx}: no existe")
    prs.save(PPTX_PATH)
    print(f"\nGuardado: {PPTX_PATH}")


if __name__ == "__main__":
    main()
