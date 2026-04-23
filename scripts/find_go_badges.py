from pptx import Presentation

for fname in ["dia_2.pptx", "dia_3.pptx"]:
    prs = Presentation(f"/mnt/d/Dev/Fabric/docs/slides/Modulo 4/{fname}")
    print("=" * 60)
    print(fname)
    print("=" * 60)
    found = False
    for idx, slide in enumerate(prs.slides):
        for s_idx, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t == "GO":
                        left_in = shape.left / 914400 if shape.left else 0
                        top_in = shape.top / 914400 if shape.top else 0
                        w_in = shape.width / 914400 if shape.width else 0
                        h_in = shape.height / 914400 if shape.height else 0
                        print(f"  Slide {idx+1} shape {s_idx} '{shape.name}': 'GO' @ ({left_in:.1f},{top_in:.1f}) size=({w_in:.1f}x{h_in:.1f})")
                        found = True
    if not found:
        print("  (ninguna coincidencia 'GO')")
    print()
