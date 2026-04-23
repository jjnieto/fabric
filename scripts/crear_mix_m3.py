"""
Crea mix.pptx combinando adopcion.pptx y dia_1.pptx del Modulo 3.
NO toca los archivos originales — genera un archivo nuevo.
Preserva formato e imagenes usando copia de XML con gestion de relationships.
"""
from pptx import Presentation
from copy import deepcopy
from pptx.oxml.ns import qn

BASE = "/mnt/d/Dev/Fabric/docs/slides/Modulo 3"
SRC_ADOPCION = f"{BASE}/adopcion.pptx"
SRC_DIA1 = f"{BASE}/dia_1.pptx"
DEST = f"{BASE}/mix.pptx"

# Plan: lista de (source_file, source_slide_idx_0based)
PLAN = [
    ("dia1",     0),   # Portada Modulo 3
    ("adopcion", 1),   # La promesa vs la realidad
    ("adopcion", 2),   # Ciclo de hype (con imagen)
    ("dia1",     1),   # Seccion: Casos que SI funcionan
    ("dia1",     2),   # Walmart Food Trust
    ("dia1",     3),   # We.Trade
    ("dia1",     4),   # HKMA eTradeConnect
    ("dia1",     5),   # Otros casos en produccion
    ("adopcion", 3),   # Seccion: 6 casos que no llegaron al mainstream
    ("adopcion", 4),   # TradeLens (detallado)
    ("adopcion", 5),   # B3i (detallado)
    ("adopcion", 6),   # MedRec (detallado)
    ("adopcion", 7),   # Voatz (detallado)
    ("adopcion", 8),   # De Beers Tracr (detallado)
    ("adopcion", 9),   # Honduras (detallado)
    ("dia1",    10),   # Tabla exitos vs fracasos
    ("adopcion",11),   # Patrones: las 5 barreras
    ("dia1",    11),   # Arbol de decision (texto)
    ("dia1",    12),   # Arbol de decision (imagen)
    ("dia1",    13),   # Treasure Hunt
    ("adopcion",12),   # Actividad en grupos (6 casos)
    ("dia1",    14),   # Debate: blockchain o BD
    ("adopcion",13),   # Debate: preguntas conjuntas
    ("adopcion",14),   # Pivote a blockchain privada
    ("adopcion",15),   # Casos que SI funcionan con privada
    ("dia1",    15),   # Repaso del dia
    # Respuestas (combinadas)
    ("adopcion",16),   # Respuestas a la actividad (1/2)
    ("adopcion",17),   # Respuestas a la actividad (2/2)
    ("dia1",    16),   # Respuestas al debate ¿blockchain o BD?
    ("adopcion",18),   # Respuestas al debate (1/2)
    ("adopcion",19),   # Respuestas al debate (2/2)
    ("dia1",    17),   # Respuestas al repaso
]


def copy_slide(src_prs, src_idx, dest_prs):
    """
    Copia una slide de src_prs a dest_prs preservando formato e imagenes.
    Tecnica: deepcopy del XML + copia de relationships (imagenes, etc).
    """
    src_slide = src_prs.slides[src_idx]

    # Usar layout blank (6) como base
    blank_layout = dest_prs.slide_layouts[6]
    new_slide = dest_prs.slides.add_slide(blank_layout)

    # Eliminar placeholders del layout
    for shape in list(new_slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    # Copiar shapes del source (deepcopy del XML)
    src_spTree = src_slide.shapes._spTree
    dest_spTree = new_slide.shapes._spTree

    # Copiar todos los hijos del spTree menos nvGrpSpPr y grpSpPr (que ya estan)
    for el in src_spTree:
        if el.tag.endswith('}nvGrpSpPr') or el.tag.endswith('}grpSpPr'):
            continue
        new_el = deepcopy(el)
        dest_spTree.append(new_el)

    # Copiar relationships (imagenes, media, etc)
    rId_map = {}
    for rId, rel in src_slide.part.rels.items():
        # rel.reltype identifica el tipo: image, chart, ole, etc
        if rel.is_external:
            # Relacion externa (hyperlinks, etc): solo copiar la URL
            new_rId = new_slide.part.relate_to(rel.target_ref, rel.reltype, is_external=True)
        else:
            # Relacion interna: copiar la parte (ej. imagen)
            target_part = rel.target_part
            new_rId = new_slide.part.relate_to(target_part, rel.reltype)
        rId_map[rId] = new_rId

    # Actualizar referencias rId en el XML copiado (r:embed, r:link, r:id)
    for el in dest_spTree.iter():
        for attr_name in list(el.attrib):
            # Atributos con rId: {http://...}embed, {http://...}link, {http://...}id
            if attr_name.endswith('}embed') or attr_name.endswith('}link') or attr_name.endswith('}id'):
                old_rId = el.attrib[attr_name]
                if old_rId in rId_map:
                    el.attrib[attr_name] = rId_map[old_rId]

    return new_slide


def main():
    # Crear presentacion destino vacia con mismo tamano que los fuentes
    src_adopcion = Presentation(SRC_ADOPCION)
    src_dia1 = Presentation(SRC_DIA1)

    dest = Presentation()
    dest.slide_width = src_adopcion.slide_width
    dest.slide_height = src_adopcion.slide_height

    # Quitar las slides por defecto si las hay
    while len(dest.slides) > 0:
        rId = dest.slides._sldIdLst[0].get(qn('r:id'))
        dest.part.drop_rel(rId)
        dest.slides._sldIdLst.remove(dest.slides._sldIdLst[0])

    sources = {"adopcion": src_adopcion, "dia1": src_dia1}

    for i, (src_name, idx) in enumerate(PLAN):
        src = sources[src_name]
        try:
            copy_slide(src, idx, dest)
            print(f"  OK [{i+1:2d}] {src_name}[{idx+1}]")
        except Exception as e:
            print(f"  FAIL [{i+1:2d}] {src_name}[{idx+1}]: {e}")

    dest.save(DEST)
    print(f"\nGuardado: {DEST}")
    print(f"Total slides: {len(dest.slides)}")


if __name__ == "__main__":
    main()
