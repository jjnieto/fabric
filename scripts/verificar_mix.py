from pptx import Presentation
try:
    prs = Presentation("/mnt/d/Dev/Fabric/docs/slides/Modulo 3/mix.pptx")
    print(f"OK: se puede abrir, {len(prs.slides)} slides")
    pics_total = 0
    for idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.shape_type == 13:
                pics_total += 1
                print(f"  Slide {idx+1}: imagen encontrada")
    print(f"Total imagenes: {pics_total}")
except Exception as e:
    print(f"ERROR: {e}")
